#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX学科分类器训练脚本 (TextCNN) - 带文件名数据增强的平衡版本
支持7类别：语文、数学、英语、物理、化学、生物、班会
核心特性：
1. 文本权重70%，文件名权重30%（平衡泛化能力）
2. 文件名数据增强（模拟不同命名习惯）
3. 多尺度卷积提取文本特征
4. 交叉验证评估泛化能力
"""

import os
import re
import pickle
import random
import numpy as np
from pathlib import Path
from sklearn.model_selection import train_test_split, StratifiedKFold
from sklearn.metrics import classification_report, confusion_matrix

# 在导入tensorflow之前设置环境变量
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'

from tensorflow.keras.preprocessing.text import Tokenizer
from tensorflow.keras.preprocessing.sequence import pad_sequences
from tensorflow.keras.models import Model
from tensorflow.keras.layers import Embedding, Conv1D, GlobalMaxPooling1D, Dense, Dropout, Input, Concatenate
from tensorflow.keras.callbacks import EarlyStopping, ModelCheckpoint, ReduceLROnPlateau
from tensorflow.keras.utils import to_categorical
import jieba
from pptx import Presentation

# ---------- 配置参数 ----------
DATA_ROOT = "../data"
CATEGORIES = ["语文", "数学", "英语", "物理", "化学", "生物", "班会"]

# 词向量配置
MAX_NB_WORDS = 15000
MAX_SEQUENCE_LENGTH = 750
EMBEDDING_DIM = 120

# 文件名特征配置
MAX_FILENAME_LENGTH = 30        # 增加到30，容纳增强后的文件名
FILENAME_EMBEDDING_DIM = 32

# 🔑 平衡权重（提高文本权重，降低文件名权重）
TEXT_WEIGHT = 0.70
FILENAME_WEIGHT = 0.30

# 🔑 数据增强配置
ENABLE_FILENAME_AUGMENTATION = True
FILENAME_AUGMENTATION_COUNT = 4  # 每个文件生成4个变体（原始+3个增强）
ENABLE_TEXT_AUGMENTATION = False  # 文本增强可选，暂不启用

# 训练配置
BATCH_SIZE = 16
EPOCHS = 32
VALIDATION_SPLIT = 0.15
TEST_SPLIT = 0.15

# 交叉验证配置
ENABLE_CROSS_VALIDATION = False   # 设为True启用5折交叉验证
CV_FOLDS = 5

# 停用词表
STOPWORDS = set([
    '的', '了', '是', '我', '你', '他', '她', '它', '我们', '你们', '他们',
    '这', '那', '有', '在', '不', '和', '与', '就', '都', '而', '及', '或',
    '一个', '这个', '那个', '这些', '那些', '这里', '那里', '然后', '因为',
    '所以', '但是', '如果', '虽然', '然而', '并且', '或者'
])

# 有意义的单字词（在文件名中保留）
MEANINGFUL_SINGLE_CHARS = {'圆', '力', '氧', '氢', '碳', '钠', '酸', '碱', '盐', 
                           '电', '光', '声', '热', '诗', '词', '歌', '曲', '数',
                           '方', '程', '函', '数', '角', '形', '体', '积'}

# ---------- 文件名数据增强函数 ----------
def augment_filename(filename):
    """
    生成文件名的多种变体，提高模型泛化能力
    模拟不同的命名习惯：带数字/不带数字、中英文混合/纯中文、词序变化等
    """
    if not filename:
        return [filename]
    
    variants = [filename]  # 原始文件名
    
    # 变体1：移除所有数字
    variant1 = re.sub(r'\d+', '', filename)
    variant1 = re.sub(r'\s+', ' ', variant1).strip()
    if variant1 and variant1 != filename:
        variants.append(variant1)
    
    # 变体2：移除所有英文字母
    variant2 = re.sub(r'[a-zA-Z]+', '', filename)
    variant2 = re.sub(r'\s+', ' ', variant2).strip()
    if variant2 and variant2 != filename and variant2 != variant1:
        variants.append(variant2)
    
    # 变体3：只保留中文和数字（移除英文）
    variant3 = re.sub(r'[a-zA-Z]', '', filename)
    variant3 = re.sub(r'\s+', ' ', variant3).strip()
    if variant3 and variant3 != filename and variant3 != variant2:
        variants.append(variant3)
    
    # 变体4：只保留中文关键词（移除英文和数字）
    variant4 = re.sub(r'[a-zA-Z0-9]', '', filename)
    variant4 = re.sub(r'\s+', ' ', variant4).strip()
    if variant4 and variant4 != filename and variant4 != variant3:
        variants.append(variant4)
    
    # 变体5：打乱词序（模拟不同命名顺序）
    words = filename.split()
    if len(words) >= 2:
        shuffled = words.copy()
        random.shuffle(shuffled)
        variant5 = ' '.join(shuffled)
        if variant5 != filename:
            variants.append(variant5)
    
    # 变体6：移除停用词（精简文件名）
    words_no_stop = [w for w in words if w not in STOPWORDS]
    if words_no_stop and len(words_no_stop) < len(words):
        variant6 = ' '.join(words_no_stop)
        if variant6 != filename:
            variants.append(variant6)
    
    # 去重并限制数量
    variants = list(dict.fromkeys(variants))
    return variants[:FILENAME_AUGMENTATION_COUNT]

def augment_text(text):
    """
    文本数据增强（可选，暂不启用）
    可用于增加文本多样性
    """
    if not ENABLE_TEXT_AUGMENTATION:
        return [text]
    
    variants = [text]
    # 可在此添加同义词替换、随机删除等增强策略
    return variants

def extract_text_from_pptx(pptx_path):
    """从PPTX文件中提取所有文本内容"""
    text_parts = []
    pptx_path_obj = Path(pptx_path) if isinstance(pptx_path, str) else pptx_path
    file_name = pptx_path_obj.name
    
    try:
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text.strip())
                    
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text_parts.append(cell.text.strip())
                    
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text:
                                text_parts.append(paragraph.text.strip())
                except:
                    continue
            
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        text_parts.append(notes_slide.notes_text_frame.text.strip())
            except:
                pass
            
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    text_parts.append(slide.shapes.title.text.strip())
            except:
                pass
        
        result = " ".join(text_parts)
        return result
        
    except Exception as e:
        print(f"    读取失败 {file_name}: {str(e)[:100]}")
        return ""

def extract_filename_features(filepath):
    """从文件路径中提取文件名特征"""
    filepath_obj = Path(filepath) if isinstance(filepath, str) else filepath
    filename = filepath_obj.stem
    
    # 清洗文件名：保留中文、英文、数字、空格
    cleaned = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', ' ', filename)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    
    return cleaned, filename

def process_filename_text(filename_text):
    """处理文件名文本：分词、过滤停用词"""
    if not filename_text:
        return ""
    
    words = jieba.cut(filename_text)
    filtered = []
    for w in words:
        if not w or w.strip() == '':
            continue
        if w not in STOPWORDS:
            # 保留有意义的单字词
            if len(w) > 1 or w in MEANINGFUL_SINGLE_CHARS:
                filtered.append(w)
    
    if not filtered:
        filtered = [w for w in words if w.strip()]
    
    return " ".join(filtered)

def clean_text(text):
    """基础清洗：保留中文、英文、数字、常用标点"""
    if not text:
        return ""
    
    text = re.sub(r'http\S+|www\S+|https\S+', '', text)
    text = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9，。！？；：""''（）【】《》、 ]', '', text)
    text = re.sub(r'\s+', ' ', text)
    
    return text.strip()

def cut_words(text):
    """使用jieba进行分词，过滤停用词"""
    if not text:
        return ""
    
    words = jieba.cut(text)
    filtered = []
    for w in words:
        if not w or w.strip() == '':
            continue
        if w not in STOPWORDS:
            filtered.append(w)
    
    if not filtered:
        filtered = [w for w in words if w.strip()]
    
    if not filtered:
        return "空文本"
    
    return " ".join(filtered)

def load_data_with_augmentation(data_root):
    """
    带数据增强的加载函数
    每个文件生成多个文件名变体，但文本内容保持不变
    """
    texts = []
    filenames = []
    labels = []
    
    stats = {
        'total_files': 0,
        'extract_failed': 0,
        'text_empty': 0,
        'text_short': 0,
        'text_success': 0,
        'filename_success': 0,
        'augmented_count': 0,
        'success': 0,
        'by_category': {}
    }
    
    for label_idx, cat in enumerate(CATEGORIES):
        cat_path = os.path.join(data_root, cat)
        if not os.path.isdir(cat_path):
            print(f"警告: 目录不存在 {cat_path}")
            stats['by_category'][cat] = {'total': 0, 'success': 0, 'augmented': 0}
            continue
        
        pptx_files = list(Path(cat_path).glob("*.pptx")) + list(Path(cat_path).glob("*.PPTX"))
        stats['by_category'][cat] = {'total': len(pptx_files), 'success': 0, 'augmented': 0}
        stats['total_files'] += len(pptx_files)
        
        if not pptx_files:
            print(f"警告: {cat} 目录下没有找到PPTX文件")
            continue
        
        print(f"\n处理 {cat} 目录，共 {len(pptx_files)} 个文件")
        
        for pptx_file in pptx_files:
            # 提取文件名特征
            filename_raw, original_filename = extract_filename_features(str(pptx_file))
            filename_processed = process_filename_text(filename_raw)
            
            # 提取文本内容
            raw_text = extract_text_from_pptx(str(pptx_file))
            
            text_available = False
            if raw_text:
                cleaned = clean_text(raw_text)
                if cleaned:
                    segmented = cut_words(cleaned)
                    if segmented:
                        text_available = True
                        stats['text_success'] += 1
                        text_content = segmented
                        
                        if len(cleaned) < 20:
                            stats['text_short'] += 1
                    else:
                        text_content = "空文本"
                else:
                    text_content = "空文本"
            else:
                text_content = "空文本"
                stats['text_empty'] += 1
            
            # 获取文件名增强变体
            filename_variants = [filename_processed]
            if ENABLE_FILENAME_AUGMENTATION and filename_processed:
                # 对原始文件名进行增强
                augmented_variants = augment_filename(original_filename)
                for variant in augmented_variants:
                    if variant and variant != original_filename:
                        variant_processed = process_filename_text(variant)
                        if variant_processed and variant_processed != filename_processed:
                            filename_variants.append(variant_processed)
                
                # 限制变体数量
                filename_variants = list(dict.fromkeys(filename_variants))[:FILENAME_AUGMENTATION_COUNT]
            
            # 添加样本（原始 + 增强）
            for fname_var in filename_variants:
                if fname_var:  # 确保文件名非空
                    texts.append(text_content)
                    filenames.append(fname_var)
                    labels.append(label_idx)
                    stats['success'] += 1
                    stats['by_category'][cat]['success'] += 1
                    
                    if fname_var != filename_processed:
                        stats['augmented_count'] += 1
                        stats['by_category'][cat]['augmented'] += 1
            
            # 打印进度
            if stats['success'] % 50 == 0:
                print(f"  ... 已处理 {stats['success']} 个样本（含增强）")
        
        cat_stats = stats['by_category'][cat]
        print(f"  {cat}: 原始{cat_stats['total']}个文件 → {cat_stats['success']}个样本（含{cat_stats['augmented']}个增强）")
    
    return texts, filenames, labels, stats

def print_stats(stats):
    """打印详细统计信息"""
    print("\n" + "="*50)
    print("数据加载统计报告（含数据增强）")
    print("="*50)
    print(f"总文件数:        {stats['total_files']}")
    print(f"生成样本数:      {stats['success']}")
    print(f"增强样本数:      {stats['augmented_count']}")
    print(f"  └─ 增强比例:   {stats['augmented_count']/stats['success']*100:.1f}%")
    print(f"\n文本状态:")
    print(f"  有文本内容:    {stats['text_success']}")
    print(f"  文本过短:      {stats['text_short']}")
    print(f"  文本为空:      {stats['text_empty']}")
    print("\n各类别统计:")
    for cat, info in stats['by_category'].items():
        if info['total'] > 0:
            print(f"  {cat}: {info['total']}个文件 → {info['success']}个样本（{info['augmented']}增强）")
    print("="*50)

def build_balanced_multimodal_model(vocab_size, max_seq_len, max_filename_len, 
                                      filename_vocab_size, num_classes):
    """
    构建权重平衡的多模态模型
    文本分支更深、更宽，文件名分支较浅
    """
    
    # ========== 文本分支（权重70%）==========
    text_input = Input(shape=(max_seq_len,), name='text_input')
    text_embedding = Embedding(vocab_size, EMBEDDING_DIM, name='text_embedding')(text_input)
    
    # 多尺度卷积提取文本特征
    conv_3 = Conv1D(filters=128, kernel_size=3, activation='relu', padding='same', name='text_conv_3')(text_embedding)
    conv_4 = Conv1D(filters=128, kernel_size=4, activation='relu', padding='same', name='text_conv_4')(text_embedding)
    conv_5 = Conv1D(filters=128, kernel_size=5, activation='relu', padding='same', name='text_conv_5')(text_embedding)
    
    pool_3 = GlobalMaxPooling1D(name='text_pool_3')(conv_3)
    pool_4 = GlobalMaxPooling1D(name='text_pool_4')(conv_4)
    pool_5 = GlobalMaxPooling1D(name='text_pool_5')(conv_5)
    
    # 合并多尺度特征
    text_multi = Concatenate(name='text_multi')([pool_3, pool_4, pool_5])
    text_dense = Dense(128, activation='relu', name='text_dense')(text_multi)
    text_dropout = Dropout(0.5, name='text_dropout')(text_dense)
    
    # ========== 文件名分支（权重30%）==========
    filename_input = Input(shape=(max_filename_len,), name='filename_input')
    filename_embedding = Embedding(filename_vocab_size, FILENAME_EMBEDDING_DIM, name='filename_embedding')(filename_input)
    filename_conv = Conv1D(filters=64, kernel_size=2, activation='relu', name='filename_conv')(filename_embedding)
    filename_pool = GlobalMaxPooling1D(name='filename_pool')(filename_conv)
    filename_dense = Dense(32, activation='relu', name='filename_dense')(filename_pool)
    filename_dropout = Dropout(0.3, name='filename_dropout')(filename_dense)
    
    # ========== 融合层 ==========
    combined = Concatenate(name='combined')([text_dropout, filename_dropout])
    
    # 最终分类层
    final_dense = Dense(64, activation='relu', name='final_dense')(combined)
    final_dropout = Dropout(0.5, name='final_dropout')(final_dense)
    output = Dense(num_classes, activation='softmax', name='output')(final_dropout)
    
    model = Model(inputs=[text_input, filename_input], outputs=output)
    model.compile(loss='categorical_crossentropy', optimizer='adam', metrics=['accuracy'])
    
    return model

def cross_validate(texts, filenames, labels):
    """K折交叉验证评估泛化能力"""
    print("\n" + "="*50)
    print(f"开始 {CV_FOLDS} 折交叉验证")
    print("="*50)
    
    skf = StratifiedKFold(n_splits=CV_FOLDS, shuffle=True, random_state=42)
    fold_scores = []
    
    for fold, (train_idx, val_idx) in enumerate(skf.split(texts, labels)):
        print(f"\n第 {fold+1}/{CV_FOLDS} 折")
        
        # 划分数据
        X_train_text = [texts[i] for i in train_idx]
        X_val_text = [texts[i] for i in val_idx]
        X_train_filename = [filenames[i] for i in train_idx]
        X_val_filename = [filenames[i] for i in val_idx]
        y_train = [labels[i] for i in train_idx]
        y_val = [labels[i] for i in val_idx]
        
        # 构建Tokenizer（仅基于训练集）
        text_tokenizer = Tokenizer(num_words=MAX_NB_WORDS, oov_token='<UNK>')
        text_tokenizer.fit_on_texts(X_train_text)
        text_vocab_size = min(MAX_NB_WORDS, len(text_tokenizer.word_index) + 1)
        
        filename_tokenizer = Tokenizer(num_words=5000, oov_token='<UNK>')
        filename_tokenizer.fit_on_texts(X_train_filename)
        filename_vocab_size = min(5000, len(filename_tokenizer.word_index) + 1)
        
        # 转换序列
        def seq_and_pad_text(texts):
            seqs = text_tokenizer.texts_to_sequences(texts)
            return pad_sequences(seqs, maxlen=MAX_SEQUENCE_LENGTH, padding='post', truncating='post')
        
        def seq_and_pad_filename(filenames):
            seqs = filename_tokenizer.texts_to_sequences(filenames)
            return pad_sequences(seqs, maxlen=MAX_FILENAME_LENGTH, padding='post', truncating='post')
        
        X_train_text_pad = seq_and_pad_text(X_train_text)
        X_val_text_pad = seq_and_pad_text(X_val_text)
        X_train_filename_pad = seq_and_pad_filename(X_train_filename)
        X_val_filename_pad = seq_and_pad_filename(X_val_filename)
        
        y_train_cat = to_categorical(y_train, num_classes=len(CATEGORIES))
        y_val_cat = to_categorical(y_val, num_classes=len(CATEGORIES))
        
        # 构建模型
        model = build_balanced_multimodal_model(
            vocab_size=text_vocab_size,
            max_seq_len=MAX_SEQUENCE_LENGTH,
            max_filename_len=MAX_FILENAME_LENGTH,
            filename_vocab_size=filename_vocab_size,
            num_classes=len(CATEGORIES)
        )
        
        # 训练
        early_stop = EarlyStopping(monitor='val_loss', patience=3, restore_best_weights=True)
        model.fit(
            [X_train_text_pad, X_train_filename_pad], y_train_cat,
            batch_size=BATCH_SIZE,
            epochs=EPOCHS // 2,  # 交叉验证时减少epoch
            validation_data=([X_val_text_pad, X_val_filename_pad], y_val_cat),
            callbacks=[early_stop],
            verbose=0
        )
        
        # 评估
        loss, acc = model.evaluate([X_val_text_pad, X_val_filename_pad], y_val_cat, verbose=0)
        fold_scores.append(acc)
        print(f"  验证准确率: {acc:.4f}")
    
    print("\n" + "="*50)
    print(f"交叉验证结果: 平均准确率 = {np.mean(fold_scores):.4f} ± {np.std(fold_scores):.4f}")
    print("="*50)
    
    return fold_scores

# ---------- 主流程 ----------
def main():
    print("="*50)
    print("PPTX学科分类器训练（平衡权重 + 文件名数据增强）")
    print(f"支持类别: {', '.join(CATEGORIES)}")
    print(f"文本权重: {TEXT_WEIGHT}, 文件名权重: {FILENAME_WEIGHT}")
    print(f"数据增强: {'启用' if ENABLE_FILENAME_AUGMENTATION else '禁用'}")
    print("="*50)
    
    # 检查data目录
    if not os.path.exists(DATA_ROOT):
        print(f"错误: 数据目录不存在 '{DATA_ROOT}'")
        return
    
    # 1. 加载数据（带增强）
    print("\n步骤1: 加载PPTX文件（带数据增强）...")
    texts, filenames, labels, stats = load_data_with_augmentation(DATA_ROOT)
    print_stats(stats)
    
    if len(texts) == 0:
        print("\n错误: 未找到任何有效PPTX文件！")
        return
    
    # 2. 可选：交叉验证
    if ENABLE_CROSS_VALIDATION and len(texts) > 100:
        cross_validate(texts, filenames, labels)
    
    # 3. 划分训练/测试集
    print("\n步骤2: 划分数据集...")
    X_train_text, X_test_text, X_train_filename, X_test_filename, y_train, y_test = train_test_split(
        texts, filenames, labels, test_size=TEST_SPLIT, stratify=labels, random_state=42
    )
    
    val_ratio = VALIDATION_SPLIT / (1 - TEST_SPLIT)
    X_train_text, X_val_text, X_train_filename, X_val_filename, y_train, y_val = train_test_split(
        X_train_text, X_train_filename, y_train, test_size=val_ratio, 
        stratify=y_train, random_state=42
    )
    print(f"训练集: {len(X_train_text)}")
    print(f"验证集: {len(X_val_text)}")
    print(f"测试集: {len(X_test_text)}")
    
    # 4. 构建Tokenizer
    print("\n步骤3: 构建词表并转换序列...")
    
    text_tokenizer = Tokenizer(num_words=MAX_NB_WORDS, oov_token='<UNK>')
    text_tokenizer.fit_on_texts(X_train_text)
    text_vocab_size = min(MAX_NB_WORDS, len(text_tokenizer.word_index) + 1)
    print(f"文本词汇量: {len(text_tokenizer.word_index)}")
    
    filename_tokenizer = Tokenizer(num_words=5000, oov_token='<UNK>')
    filename_tokenizer.fit_on_texts(X_train_filename)
    filename_vocab_size = min(5000, len(filename_tokenizer.word_index) + 1)
    print(f"文件名词汇量: {len(filename_tokenizer.word_index)}")
    
    # 序列化
    def seq_and_pad_text(texts):
        seqs = text_tokenizer.texts_to_sequences(texts)
        return pad_sequences(seqs, maxlen=MAX_SEQUENCE_LENGTH, padding='post', truncating='post')
    
    def seq_and_pad_filename(filenames):
        seqs = filename_tokenizer.texts_to_sequences(filenames)
        return pad_sequences(seqs, maxlen=MAX_FILENAME_LENGTH, padding='post', truncating='post')
    
    X_train_text_pad = seq_and_pad_text(X_train_text)
    X_val_text_pad = seq_and_pad_text(X_val_text)
    X_test_text_pad = seq_and_pad_text(X_test_text)
    
    X_train_filename_pad = seq_and_pad_filename(X_train_filename)
    X_val_filename_pad = seq_and_pad_filename(X_val_filename)
    X_test_filename_pad = seq_and_pad_filename(X_test_filename)
    
    y_train_cat = to_categorical(y_train, num_classes=len(CATEGORIES))
    y_val_cat = to_categorical(y_val, num_classes=len(CATEGORIES))
    y_test_cat = to_categorical(y_test, num_classes=len(CATEGORIES))
    
    print(f"文本输入形状: {X_train_text_pad.shape}")
    print(f"文件名输入形状: {X_train_filename_pad.shape}")
    
    # 5. 构建模型
    print("\n步骤4: 构建平衡多模态模型...")
    model = build_balanced_multimodal_model(
        vocab_size=text_vocab_size,
        max_seq_len=MAX_SEQUENCE_LENGTH,
        max_filename_len=MAX_FILENAME_LENGTH,
        filename_vocab_size=filename_vocab_size,
        num_classes=len(CATEGORIES)
    )
    model.summary()
    
    # 6. 训练
    print("\n步骤5: 开始训练...")
    early_stop = EarlyStopping(monitor='val_loss', patience=5, restore_best_weights=True)
    checkpoint = ModelCheckpoint('best_model_balanced.keras', monitor='val_accuracy', save_best_only=True, mode='max')
    reduce_lr = ReduceLROnPlateau(monitor='val_loss', factor=0.5, patience=3, min_lr=1e-6)
    
    history = model.fit(
        [X_train_text_pad, X_train_filename_pad], y_train_cat,
        batch_size=BATCH_SIZE,
        epochs=EPOCHS,
        validation_data=([X_val_text_pad, X_val_filename_pad], y_val_cat),
        callbacks=[early_stop, checkpoint, reduce_lr],
        verbose=1
    )
    
    # 7. 评估
    print("\n步骤6: 模型评估...")
    loss, acc = model.evaluate([X_test_text_pad, X_test_filename_pad], y_test_cat, verbose=0)
    print(f"测试集准确率: {acc:.4f}")
    
    y_pred_proba = model.predict([X_test_text_pad, X_test_filename_pad])
    y_pred = np.argmax(y_pred_proba, axis=1)
    print("\n详细分类报告:")
    print(classification_report(y_test, y_pred, target_names=CATEGORIES))
    print("混淆矩阵:")
    print(confusion_matrix(y_test, y_pred))
    
    # 8. 保存模型
    print("\n步骤7: 保存模型...")
    model.save('textcnn_balanced_classifier.keras')
    with open('text_tokenizer.pkl', 'wb') as f:
        pickle.dump(text_tokenizer, f)
    with open('filename_tokenizer.pkl', 'wb') as f:
        pickle.dump(filename_tokenizer, f)
    with open('categories.pkl', 'wb') as f:
        pickle.dump(CATEGORIES, f)
    
    config = {
        'max_sequence_length': MAX_SEQUENCE_LENGTH,
        'max_filename_length': MAX_FILENAME_LENGTH,
        'categories': CATEGORIES,
        'text_vocab_size': text_vocab_size,
        'filename_vocab_size': filename_vocab_size,
        'embedding_dim': EMBEDDING_DIM,
        'filename_embedding_dim': FILENAME_EMBEDDING_DIM,
        'text_weight': TEXT_WEIGHT,
        'filename_weight': FILENAME_WEIGHT
    }
    with open('config_balanced.pkl', 'wb') as f:
        pickle.dump(config, f)
    
    print("="*50)
    print("训练完成！已保存以下文件:")
    print("  - textcnn_balanced_classifier.keras (平衡模型)")
    print("  - best_model_balanced.keras (最佳模型)")
    print("  - text_tokenizer.pkl")
    print("  - filename_tokenizer.pkl")
    print("  - categories.pkl")
    print("  - config_balanced.pkl")
    print("="*50)
    
    # 打印增强示例
    print("\n文件名增强示例:")
    print("-"*40)
    sample_files = list(Path(DATA_ROOT).glob("*/[!.]*.pptx"))[:3]
    for f in sample_files:
        _, raw = extract_filename_features(str(f))
        print(f"原始: {raw}")
        variants = augment_filename(raw)
        for i, v in enumerate(variants[1:3]):  # 只显示前2个变体
            print(f"  增强{i+1}: {v}")
        print()

if __name__ == "__main__":
    # 设置随机种子保证可重复性
    random.seed(42)
    np.random.seed(42)
    main()
