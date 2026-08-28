#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX/PPT/DOCX学科分类器训练脚本 (TextCNN) - 优化版（缓存修复版）
支持8类别：语文、数学、英语、物理、化学、生物、班会
支持文件格式：.pptx, .ppt, .docx
核心特性：
1. 文本权重70%，文件名权重30%（增强文本重要性）
2. 文件名数据增强（模拟不同命名习惯）
3. 多尺度卷积提取文本特征（增加kernel_size=2）
4. 🔥 文件解析缓存（立即保存，避免重复解析，大幅加速）
5. 类别权重处理（解决小样本类别不平衡）
6. 🎯 支持指定科目只使用缓存（跳过解析，仅从缓存读取）
7. 👤 人名去除（在缓存阶段使用jieba识别并过滤人名）
8. 📊 词表截断（限制在MAX_NB_WORDS，如20000词）
9. 💾 自动保存无依赖的Tokenizer（用于推理）
10. 📝 输出JSON格式词表（便于外部工具使用）
11. 📄 支持DOCX文件解析
12. 🔍 增强缓存调试和验证
"""

import os
import re
import pickle
import json
import random
import zipfile
from io import BytesIO
import numpy as np
from pathlib import Path
from datetime import datetime
from sklearn.model_selection import train_test_split
from sklearn.metrics import classification_report, confusion_matrix
from sklearn.utils.class_weight import compute_class_weight

# 在导入tensorflow之前设置环境变量
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'

from tensorflow.keras.preprocessing.text import Tokenizer
from tensorflow.keras.preprocessing.sequence import pad_sequences
from tensorflow.keras.models import Model
from tensorflow.keras.layers import Embedding, Conv1D, GlobalMaxPooling1D, Dense, Dropout, Input, Concatenate
from tensorflow.keras.callbacks import EarlyStopping, ModelCheckpoint, ReduceLROnPlateau
from tensorflow.keras.utils import to_categorical
import jieba
import jieba.posseg as pseg
from pptx import Presentation

# 尝试导入docx处理库
try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("警告: python-docx未安装，docx文件支持将受限。请运行: pip install python-docx")

# ========== 优化后的配置参数 ==========
DATA_ROOT = "../data"
CACHE_DIR = "../cache"
CATEGORIES = ["语文", "数学", "英语", "物理", "化学", "生物", "班会"]

# 🎯 指定只使用缓存的科目（跳过解析，仅从已有缓存读取）
SKIP_CATEGORIES = []  # 例如: ["语文", "数学"] 表示语文和数学只从缓存读取

# 词向量配置（增强语义表示）
MAX_NB_WORDS = 20000          # 词表最大大小（在20000词处截断）
MAX_SEQUENCE_LENGTH = 1000    # 增加序列长度
EMBEDDING_DIM = 150           # 增强Embedding维度

# 文件名特征配置
MAX_FILENAME_LENGTH = 32
FILENAME_EMBEDDING_DIM = 32
MAX_FILENAME_WORDS = 5000     # 文件名词表最大大小

# 平衡权重（增强文本重要性）
TEXT_WEIGHT = 0.70            # 提高文本权重
FILENAME_WEIGHT = 0.30        # 降低文件名权重

# 数据增强配置（减少噪声）
ENABLE_FILENAME_AUGMENTATION = True
FILENAME_AUGMENTATION_COUNT = 3  # 减少变体数量

# 训练配置
BATCH_SIZE = 12               # 更精细的梯度更新
EPOCHS = 40                   # 更多训练轮次
VALIDATION_SPLIT = 0.15
TEST_SPLIT = 0.15

# 交叉验证配置
ENABLE_CROSS_VALIDATION = False
CV_FOLDS = 5

# 缓存配置
ENABLE_CACHE = True
CACHE_VERSION = "v6"          # 更新缓存版本（修复保存问题）
FORCE_REFRESH_CACHE = False

# 类别权重配置
USE_CLASS_WEIGHTS = True      # 启用类别权重（处理小样本）

# 人名去除配置
REMOVE_PERSON_NAMES = True    # 是否去除人名
KEEP_SINGLE_CHAR_NAMES = False # 是否保留单字人名（通常单字可能是误判）

# 支持的文档格式
SUPPORTED_EXTENSIONS = ['.pptx', '.ppt', '.docx', '.DOCX', '.PPTX', '.PPT']

# 停用词表
STOPWORDS = set([
    '的', '了', '是', '我', '你', '他', '她', '它', '我们', '你们', '他们',
    '这', '那', '有', '在', '不', '和', '与', '就', '都', '而', '及', '或',
    '一个', '这个', '那个', '这些', '那些', '这里', '那里', '然后', '因为',
    '所以', '但是', '如果', '虽然', '然而', '并且', '或者'
])

# 有意义的单字词
MEANINGFUL_SINGLE_CHARS = {'圆', '力', '氧', '氢', '碳', '钠', '酸', '碱', '盐', 
                           '电', '光', '声', '热', '诗', '词', '歌', '曲', '数',
                           '方', '程', '函', '数', '角', '形', '体', '积'}

# 常见姓氏（用于人名识别增强）
COMMON_SURNAMES = set([
    '李', '王', '张', '刘', '陈', '杨', '赵', '黄', '周', '吴', '徐', '孙', '马', '朱',
    '胡', '林', '郭', '何', '高', '郑', '罗', '梁', '谢', '宋', '唐', '邓', '萧', '冯',
    '韩', '曹', '彭', '曾', '肖', '田', '董', '潘', '袁', '于', '蒋', '蔡', '余', '杜',
    '苏', '吕', '丁', '沈', '任', '姚', '卢', '傅', '钟', '姜', '崔', '谭', '廖', '范',
    '汪', '陆', '金', '石', '戴', '贾', '韦', '夏', '邱', '方', '侯', '邹', '熊', '孟',
    '秦', '白', '江', '阎', '薛', '尹', '段', '雷', '黎', '史', '龙', '陶', '贺', '顾',
    '毛', '郝', '龚', '邵', '万', '钱', '严', '赖', '覃', '洪', '武', '莫', '孔', '汤',
    '向', '常', '温', '康', '施', '文', '牛', '樊', '葛', '邢'
])

# ---------- 人名识别和去除模块 ----------
class PersonNameRemover:
    """使用jieba词性标注识别并去除人名"""
    
    def __init__(self):
        # 加载jieba词性标注功能
        # 人名词性标签: nr, nr1, nr2, nrj, nrf, nrs, ns, nt, nz 中的人名相关
        self.person_tags = {'nr', 'nr1', 'nr2', 'nrj', 'nrf', 'nrs'}
        self._compile_name_patterns()
    
    def _compile_name_patterns(self):
        """编译人名相关的正则模式"""
        # 匹配常见姓氏开头的2-4字中文词
        surnames_pattern = '|'.join(re.escape(s) for s in COMMON_SURNAMES)
        self.surname_pattern = re.compile(f'[{surnames_pattern}][\\u4e00-\\u9fa5]{{1,3}}')
        
        # 匹配"XX老师"、"XX同学"等称呼模式
        self.title_pattern = re.compile(r'([\u4e00-\u9fa5]{2,4})(?:老师|同学|教授|博士|校长|主任|院长|书记)')
        
        # 匹配"XX说"、"XX认为"等动作模式
        self.action_pattern = re.compile(r'([\u4e00-\u9fa5]{2,4})(?:说|道|问|答|认为|表示|指出|强调|建议|提出)')
    
    def remove_names_by_pos(self, text):
        """
        使用jieba词性标注去除人名
        返回去除人名后的文本
        """
        if not text or not REMOVE_PERSON_NAMES:
            return text
        
        try:
            words = pseg.cut(text)
            filtered_words = []
            
            for word, flag in words:
                # 判断是否为人名
                is_person = flag in self.person_tags
                
                # 对于单字词，如果不是有意义单字且被标记为人名，则过滤
                if is_person:
                    if len(word) == 1 and KEEP_SINGLE_CHAR_NAMES:
                        # 检查是否为有意义单字
                        if word in MEANINGFUL_SINGLE_CHARS:
                            filtered_words.append(word)
                    # 否则跳过（不添加）
                else:
                    filtered_words.append(word)
            
            return ''.join(filtered_words)
        except Exception as e:
            # 如果词性标注失败，返回原文本
            print(f"    警告: 词性标注人名去除失败 ({e})")
            return text
    
    def remove_names_by_regex(self, text):
        """
        使用正则表达式辅助去除可能的人名
        作为词性标注的补充
        """
        if not text:
            return text
        
        result = text
        
        # 替换"XX老师"、"XX同学"等
        result = self.title_pattern.sub(r'\1', result)
        
        # 替换"XX说"、"XX认为"等
        result = self.action_pattern.sub(r'\1', result)
        
        # 移除单独的姓氏（但保留有意义的单字）
        def replace_surname(match):
            surname = match.group(0)
            if surname not in MEANINGFUL_SINGLE_CHARS:
                return ''
            return surname
        
        return result
    
    def remove_all_names(self, text, use_regex=True):
        """
        综合使用多种方法去除人名
        优先使用词性标注，正则作为补充
        """
        if not text or not REMOVE_PERSON_NAMES:
            return text
        
        # 先使用词性标注去除
        result = self.remove_names_by_pos(text)
        
        # 可选：使用正则表达式进一步清理
        if use_regex:
            result = self.remove_names_by_regex(result)
        
        # 清理多余空格
        result = re.sub(r'\s+', ' ', result)
        
        return result.strip()


# 全局人名去除器实例
_name_remover = None

def get_name_remover():
    """获取全局人名去除器实例"""
    global _name_remover
    if _name_remover is None:
        _name_remover = PersonNameRemover()
    return _name_remover


def remove_person_names_from_text(text):
    """
    从文本中去除人名（缓存阶段使用）
    """
    if not REMOVE_PERSON_NAMES or not text:
        return text
    
    remover = get_name_remover()
    return remover.remove_all_names(text)


# ---------- DOCX 文本提取函数 ----------
def extract_text_from_docx(docx_path):
    """从DOCX文件中提取所有文本内容"""
    text_parts = []
    
    if not DOCX_SUPPORT:
        print(f"    警告: python-docx未安装，无法解析 {docx_path}")
        return ""
    
    try:
        doc = Document(docx_path)
        
        # 提取段落文本
        for paragraph in doc.paragraphs:
            if paragraph.text and paragraph.text.strip():
                text_parts.append(paragraph.text.strip())
        
        # 提取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        text_parts.append(cell.text.strip())
        
        # 提取页眉页脚
        if hasattr(doc, 'sections'):
            for section in doc.sections:
                # 页眉
                if section.header and section.header.paragraphs:
                    for para in section.header.paragraphs:
                        if para.text and para.text.strip():
                            text_parts.append(para.text.strip())
                # 页脚
                if section.footer and section.footer.paragraphs:
                    for para in section.footer.paragraphs:
                        if para.text and para.text.strip():
                            text_parts.append(para.text.strip())
        
        # 提取脚注和尾注
        if hasattr(doc, 'footnotes') and doc.footnotes:
            for footnote in doc.footnotes:
                if hasattr(footnote, 'text') and footnote.text:
                    text_parts.append(footnote.text.strip())
        
        if hasattr(doc, 'endnotes') and doc.endnotes:
            for endnote in doc.endnotes:
                if hasattr(endnote, 'text') and endnote.text:
                    text_parts.append(endnote.text.strip())
        
        result = " ".join(text_parts)
        return result
        
    except Exception as e:
        print(f"    读取DOCX失败 {Path(docx_path).name}: {str(e)[:100]}")
        return ""


def extract_text_from_docx_simple(docx_path):
    """
    简单提取DOCX文本（不使用python-docx，仅通过zip解析）
    作为fallback方案
    """
    text_parts = []
    
    try:
        with zipfile.ZipFile(docx_path, 'r') as zf:
            if 'word/document.xml' in zf.namelist():
                xml_content = zf.read('word/document.xml').decode('utf-8', errors='ignore')
                
                # 简单正则提取文本
                text_matches = re.findall(r'>([^<]+)<', xml_content)
                for match in text_matches:
                    if match.strip() and len(match.strip()) > 1:
                        text_parts.append(match.strip())
                
                # 提取表格中的文本
                # 处理 <w:t> 标签
                w_t_matches = re.findall(r'<w:t[^>]*>([^<]+)</w:t>', xml_content)
                for match in w_t_matches:
                    if match.strip():
                        text_parts.append(match.strip())
    
    except Exception as e:
        return ""
    
    return " ".join(text_parts)


# ---------- PPTX 嵌入DOCX提取 ----------
def extract_embedded_docx_from_pptx(pptx_path):
    """从PPTX文件中提取嵌入的DOCX文件并读取其内容"""
    embedded_texts = []
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            # 查找嵌入的对象
            for file_info in zf.filelist:
                if file_info.filename.startswith('ppt/embeddings/') and file_info.filename.endswith(('.docx', '.doc')):
                    try:
                        docx_data = zf.read(file_info.filename)
                        
                        if DOCX_SUPPORT:
                            docx_stream = BytesIO(docx_data)
                            doc = Document(docx_stream)
                            
                            for paragraph in doc.paragraphs:
                                if paragraph.text and paragraph.text.strip():
                                    embedded_texts.append(paragraph.text.strip())
                            
                            for table in doc.tables:
                                for row in table.rows:
                                    for cell in row.cells:
                                        if cell.text and cell.text.strip():
                                            embedded_texts.append(cell.text.strip())
                            
                            del doc
                        else:
                            # 降级方案：直接解析XML
                            try:
                                docx_zip = zipfile.ZipFile(BytesIO(docx_data))
                                if 'word/document.xml' in docx_zip.namelist():
                                    xml_content = docx_zip.read('word/document.xml').decode('utf-8', errors='ignore')
                                    text_matches = re.findall(r'>([^<]+)<', xml_content)
                                    for match in text_matches:
                                        if match.strip() and len(match.strip()) > 1:
                                            embedded_texts.append(match.strip())
                                docx_zip.close()
                            except:
                                pass
                    except Exception as e:
                        pass
    
    except Exception as e:
        pass
    
    return " ".join(embedded_texts) if embedded_texts else ""


# ---------- 统一文件解析函数 ----------
def extract_text_from_file(filepath, extract_embedded=True):
    """
    根据文件类型提取文本内容
    支持: .pptx, .ppt, .docx
    """
    filepath = Path(filepath)
    suffix = filepath.suffix.lower()
    
    if suffix in ['.pptx', '.ppt']:
        return extract_text_from_pptx(filepath, extract_embedded)
    elif suffix == '.docx':
        if DOCX_SUPPORT:
            return extract_text_from_docx(filepath)
        else:
            return extract_text_from_docx_simple(filepath)
    else:
        return ""


def extract_text_from_pptx(pptx_path, extract_embedded=True):
    """从PPTX文件中提取所有文本内容（包括嵌入的DOCX）"""
    text_parts = []
    embedded_text = ""
    
    try:
        # 提取嵌入的docx文件内容
        if extract_embedded:
            embedded_text = extract_embedded_docx_from_pptx(pptx_path)
            if embedded_text:
                text_parts.append(embedded_text)
        
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            # 标题
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    text_parts.append(slide.shapes.title.text.strip())
            except:
                pass
            
            # 形状中的文本
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text.strip())
                    
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text_parts.append(cell.text.strip())
                    
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text:
                                text_parts.append(paragraph.text.strip())
                except:
                    continue
            
            # 备注
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        text_parts.append(notes_slide.notes_text_frame.text.strip())
            except:
                pass
        
        result = " ".join(text_parts)
        return result
        
    except Exception as e:
        print(f"    读取PPTX失败 {Path(pptx_path).name}: {str(e)[:100]}")
        return ""


# ---------- 缓存管理类（增强版） ----------
class FileCache:
    """文件解析缓存管理器（支持立即保存、增强调试）"""
    
    def __init__(self, cache_dir=CACHE_DIR, version=CACHE_VERSION):
        self.cache_dir = Path(cache_dir)
        self.version = version
        self.text_cache_file = self.cache_dir / f"file_cache_{version}.json"
        self.metadata_file = self.cache_dir / f"cache_metadata_{version}.json"
        self.debug_file = self.cache_dir / f"cache_debug_{version}.log"
        
        # 创建缓存目录
        try:
            self.cache_dir.mkdir(parents=True, exist_ok=True)
            print(f"  📁 缓存目录: {self.cache_dir.absolute()}")
        except Exception as e:
            print(f"  ⚠️ 警告: 无法创建缓存目录 ({e})")
        
        # 加载现有缓存
        self.text_cache = self._load_cache()
        self.metadata = self._load_metadata()
        self._last_was_cache_hit = False
        
        # 统计信息
        self.names_removed_count = 0
        self.name_removal_stats = {}
        self.cache_save_count = 0
        self.cache_hit_count = 0
        self.cache_miss_count = 0
        
        # 调试信息
        self._debug_enabled = True
        self._debug_log = []
        
        # 打印缓存状态
        print(f"  📊 加载缓存: {len(self.text_cache)} 个文件")
        non_empty = sum(1 for v in self.text_cache.values() if v)
        if non_empty > 0:
            print(f"  ✅ 非空条目: {non_empty}")
        else:
            print(f"  ⚠️  缓存为空或全是空内容")
    
    def _log_debug(self, msg):
        """记录调试信息"""
        if self._debug_enabled:
            timestamp = datetime.now().strftime("%H:%M:%S.%f")[:-3]
            log_msg = f"[{timestamp}] {msg}"
            self._debug_log.append(log_msg)
            print(f"  🔍 {log_msg}")
    
    def _save_debug_log(self):
        """保存调试日志到文件"""
        try:
            with open(self.debug_file, 'w', encoding='utf-8') as f:
                f.write(f"缓存调试日志 - {datetime.now().isoformat()}\n")
                f.write("="*50 + "\n")
                for line in self._debug_log:
                    f.write(line + "\n")
            return True
        except Exception as e:
            return False
    
    def _get_file_hash(self, filepath):
        """计算文件的哈希值（用于检测文件是否变化）"""
        filepath = Path(filepath)
        if not filepath.exists():
            return None
        
        stat = filepath.stat()
        # 使用修改时间和文件大小作为简单hash
        return f"{stat.st_mtime}_{stat.st_size}"
    
    def _load_cache(self):
        """加载文本缓存"""
        if not ENABLE_CACHE or FORCE_REFRESH_CACHE:
            return {}
        
        if self.text_cache_file.exists():
            try:
                with open(self.text_cache_file, 'r', encoding='utf-8') as f:
                    cache = json.load(f)
                return cache
            except json.JSONDecodeError as e:
                print(f"  ⚠️  警告: 缓存文件损坏 ({e})，将重新生成")
                # 备份损坏的缓存
                backup_file = self.text_cache_file.with_suffix('.json.bak')
                try:
                    self.text_cache_file.rename(backup_file)
                    print(f"  📦 已备份损坏的缓存到: {backup_file}")
                except:
                    pass
                return {}
            except Exception as e:
                print(f"  ⚠️  警告: 加载缓存失败 ({e})，将重新生成")
                return {}
        return {}
    
    def _load_metadata(self):
        """加载缓存元数据"""
        if not ENABLE_CACHE or FORCE_REFRESH_CACHE:
            return {}
        
        if self.metadata_file.exists():
            try:
                with open(self.metadata_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except:
                return {}
        return {}
    
    def _save_cache(self):
        """保存缓存到磁盘（立即保存）"""
        if not ENABLE_CACHE:
            return
        
        try:
            # 检查是否有内容
            if not self.text_cache:
                self._log_debug("⚠️  缓存为空，跳过保存")
                return
            
            # 统计有效内容
            total_entries = len(self.text_cache)
            non_empty = sum(1 for v in self.text_cache.values() if v)
            empty_entries = total_entries - non_empty
            
            self._log_debug(f"💾 准备保存: {total_entries} 个文件，{non_empty} 个非空，{empty_entries} 个空")
            
            if non_empty == 0:
                self._log_debug("⚠️  没有非空内容，跳过保存")
                return
            
            # 保存主缓存
            with open(self.text_cache_file, 'w', encoding='utf-8') as f:
                json.dump(self.text_cache, f, ensure_ascii=False, indent=2)
            
            # 保存元数据
            with open(self.metadata_file, 'w', encoding='utf-8') as f:
                json.dump(self.metadata, f, ensure_ascii=False, indent=2)
            
            self.cache_save_count += 1
            
            # 验证保存是否成功
            if self.text_cache_file.exists():
                file_size = self.text_cache_file.stat().st_size
                self._log_debug(f"✅ 缓存已保存: {file_size} bytes")
                
                # 验证文件内容
                try:
                    with open(self.text_cache_file, 'r', encoding='utf-8') as f:
                        verify_data = json.load(f)
                        verify_count = len(verify_data)
                        if verify_count == total_entries:
                            self._log_debug(f"✅ 验证通过: {verify_count} 个条目")
                        else:
                            self._log_debug(f"⚠️  验证失败: 保存了 {verify_count} 个条目，期望 {total_entries}")
                except Exception as e:
                    self._log_debug(f"⚠️  验证失败: {e}")
            else:
                self._log_debug(f"❌ 缓存文件未创建!")
            
            # 保存人名去除统计
            if REMOVE_PERSON_NAMES and (self.names_removed_count > 0 or self.name_removal_stats):
                stats_file = self.cache_dir / f"name_removal_stats_{self.version}.json"
                with open(stats_file, 'w', encoding='utf-8') as f:
                    json.dump({
                        'names_removed_count': self.names_removed_count,
                        'name_removal_stats': self.name_removal_stats,
                        'removal_enabled': REMOVE_PERSON_NAMES,
                        'cache_version': self.version,
                        'total_files': len(self.text_cache)
                    }, f, ensure_ascii=False, indent=2)
                self._log_debug(f"✅ 人名去除统计已保存")
                
            # 保存调试日志
            self._save_debug_log()
                
        except Exception as e:
            self._log_debug(f"❌ 保存缓存失败: {e}")
            import traceback
            self._log_debug(traceback.format_exc())
    
    def get(self, filepath):
        """获取缓存的文件内容"""
        if not ENABLE_CACHE:
            self._last_was_cache_hit = False
            return None
        
        filepath = str(filepath)
        current_hash = self._get_file_hash(filepath)
        
        self._log_debug(f"📖 get: {Path(filepath).name}")
        self._log_debug(f"   - 在缓存中: {filepath in self.text_cache}")
        
        if filepath in self.text_cache:
            cached_hash = self.metadata.get(filepath, {}).get('hash')
            content = self.text_cache[filepath]
            content_len = len(content) if content else 0
            
            self._log_debug(f"   - 内容长度: {content_len}")
            self._log_debug(f"   - 哈希匹配: {cached_hash == current_hash}")
            
            if cached_hash == current_hash:
                self._last_was_cache_hit = True
                self.cache_hit_count += 1
                self._log_debug(f"   ✅ 缓存命中")
                return content
            else:
                self._log_debug(f"   ⚠️  哈希不匹配，需要重新解析")
        else:
            self._log_debug(f"   ⚠️  不在缓存中")
        
        self._last_was_cache_hit = False
        self.cache_miss_count += 1
        return None
    
    def set(self, filepath, text_content, force_save=False):
        """
        设置缓存并立即保存
        
        Args:
            filepath: 文件路径
            text_content: 文本内容
            force_save: 是否强制保存（即使内容为空）
        """
        if not ENABLE_CACHE:
            return
        
        filepath = str(filepath)
        original_len = len(text_content) if text_content else 0
        
        self._log_debug(f"💾 set: {Path(filepath).name}")
        self._log_debug(f"   - 原始内容长度: {original_len}")
        
        # 在缓存阶段应用人名去除
        if REMOVE_PERSON_NAMES and text_content:
            cleaned_content = remove_person_names_from_text(text_content)
            cleaned_len = len(cleaned_content) if cleaned_content else 0
            removed_count = original_len - cleaned_len
            
            self._log_debug(f"   - 人名去除后长度: {cleaned_len}")
            if removed_count > 0:
                self._log_debug(f"   - 去除字符数: {removed_count}")
            
            if removed_count > 0:
                self.names_removed_count += removed_count
                
                # 统计各文件的人名去除情况
                filename = Path(filepath).name
                self.name_removal_stats[filename] = {
                    'original_length': original_len,
                    'cleaned_length': cleaned_len,
                    'removed_chars': removed_count,
                    'removal_ratio': removed_count / original_len if original_len > 0 else 0
                }
            
            text_content = cleaned_content
        elif not text_content and not force_save:
            self._log_debug(f"   ⚠️  内容为空且未强制保存，跳过缓存")
            return
        
        # 更新缓存
        self.text_cache[filepath] = text_content or ""
        self.metadata[filepath] = {
            'hash': self._get_file_hash(filepath),
            'cached_at': datetime.now().isoformat(),
            'file_size': Path(filepath).stat().st_size if Path(filepath).exists() else 0,
            'names_removed': REMOVE_PERSON_NAMES,
            'file_type': Path(filepath).suffix.lower(),
            'content_length': len(text_content) if text_content else 0
        }
        
        self._log_debug(f"   ✅ 缓存已更新")
        
        # 立即保存到磁盘
        self._save_cache()
    
    def clear(self):
        """清空缓存"""
        self.text_cache = {}
        self.metadata = {}
        self.names_removed_count = 0
        self.name_removal_stats = {}
        self.cache_save_count = 0
        
        if self.text_cache_file.exists():
            self.text_cache_file.unlink()
        if self.metadata_file.exists():
            self.metadata_file.unlink()
        
        self._log_debug("🗑️  缓存已清空")
        print("  ✅ 缓存已清空")
    
    def save(self):
        """保存缓存（兼容旧接口）"""
        self._save_cache()
    
    def get_stats(self):
        """获取缓存统计信息"""
        total_entries = len(self.text_cache)
        non_empty = sum(1 for v in self.text_cache.values() if v)
        
        return {
            'cached_files': total_entries,
            'non_empty': non_empty,
            'empty_entries': total_entries - non_empty,
            'cache_file_size': self.text_cache_file.stat().st_size if self.text_cache_file.exists() else 0,
            'metadata_file_size': self.metadata_file.stat().st_size if self.metadata_file.exists() else 0,
            'names_removed_total': self.names_removed_count,
            'files_with_names_removed': len(self.name_removal_stats),
            'cache_save_count': self.cache_save_count,
            'cache_hit_count': self.cache_hit_count,
            'cache_miss_count': self.cache_miss_count
        }
    
    def has_cache(self, filepath):
        """检查文件是否有有效缓存"""
        if not ENABLE_CACHE:
            return False
        
        filepath = str(filepath)
        current_hash = self._get_file_hash(filepath)
        
        if filepath in self.text_cache:
            cached_hash = self.metadata.get(filepath, {}).get('hash')
            return cached_hash == current_hash
        return False
    
    def print_name_removal_summary(self):
        """打印人名去除摘要"""
        if not REMOVE_PERSON_NAMES:
            print("\n👤 人名去除: 已禁用")
            return
        
        print(f"\n👤 人名去除统计:")
        print(f"  总共去除字符数: {self.names_removed_count}")
        print(f"  涉及文件数: {len(self.name_removal_stats)}")
        
        # 显示去除最多的5个文件
        if self.name_removal_stats:
            sorted_files = sorted(
                self.name_removal_stats.items(),
                key=lambda x: x[1]['removed_chars'],
                reverse=True
            )[:5]
            
            if sorted_files:
                print(f"  去除最多人名的文件:")
                for filename, stats in sorted_files:
                    print(f"    - {filename[:40]}: 去除 {stats['removed_chars']} 字符 ({stats['removal_ratio']*100:.1f}%)")
    
    def print_cache_summary(self):
        """打印缓存摘要"""
        stats = self.get_stats()
        print(f"\n📊 缓存摘要:")
        print(f"  总缓存条目: {stats['cached_files']}")
        print(f"  非空条目: {stats['non_empty']}")
        print(f"  空条目: {stats['empty_entries']}")
        print(f"  缓存命中: {stats['cache_hit_count']}")
        print(f"  缓存未命中: {stats['cache_miss_count']}")
        print(f"  保存次数: {stats['cache_save_count']}")


# 全局缓存实例
_file_cache = None

def get_cache():
    """获取全局缓存实例"""
    global _file_cache
    if _file_cache is None:
        _file_cache = FileCache()
    return _file_cache


# ---------- 文件名数据增强函数 ----------
def augment_filename(filename):
    """生成文件名的多种变体（减少变体数量）"""
    if not filename:
        return [filename]
    
    variants = [filename]
    
    # 变体1：移除所有数字
    variant1 = re.sub(r'\d+', '', filename)
    variant1 = re.sub(r'\s+', ' ', variant1).strip()
    if variant1 and variant1 != filename:
        variants.append(variant1)
    
    # 变体2：只保留中文（最重要的变体）
    variant2 = re.sub(r'[a-zA-Z0-9]', '', filename)
    variant2 = re.sub(r'\s+', ' ', variant2).strip()
    if variant2 and variant2 != filename and variant2 != variant1:
        variants.append(variant2)
    
    # 去重
    variants = list(dict.fromkeys(variants))
    return variants[:FILENAME_AUGMENTATION_COUNT]


def process_filename_text(filename_text):
    """处理文件名文本：分词、过滤停用词"""
    if not filename_text:
        return ""
    
    words = jieba.cut(filename_text)
    filtered = []
    for w in words:
        if not w or w.strip() == '':
            continue
        if w not in STOPWORDS:
            if len(w) > 1 or w in MEANINGFUL_SINGLE_CHARS:
                filtered.append(w)
    
    if not filtered:
        filtered = [w for w in words if w.strip()]
    
    return " ".join(filtered)


def clean_text(text):
    """基础清洗：保留中文、英文、数字、常用标点"""
    if not text:
        return ""
    
    text = re.sub(r'http\S+|www\S+|https\S+', '', text)
    text = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9，。！？；：""''（）【】《》、 ]', '', text)
    text = re.sub(r'\s+', ' ', text)
    
    return text.strip()


def cut_words(text):
    """使用jieba进行分词，过滤停用词"""
    if not text:
        return ""
    
    words = jieba.cut(text)
    filtered = []
    for w in words:
        if not w or w.strip() == '':
            continue
        if w not in STOPWORDS:
            filtered.append(w)
    
    if not filtered:
        filtered = [w for w in words if w.strip()]
    
    if not filtered:
        return "空文本"
    
    return " ".join(filtered)


def process_file(filepath, cat, label_idx, stats, skip_parsing=False):
    """处理单个文件（支持PPTX、DOCX）- 立即缓存"""
    cache = get_cache()
    filepath_str = str(filepath)
    
    filename_raw, original_filename = extract_filename_features(filepath_str)
    filename_processed = process_filename_text(filename_raw)
    
    # 检查缓存
    cached_text = cache.get(filepath_str)
    
    if cached_text is not None and not skip_parsing:
        # 缓存命中，直接使用缓存内容
        raw_text = cached_text
        stats['cache_hits'] += 1
        stats['cache_used'] += 1
    elif skip_parsing:
        # 跳过解析模式，必须从缓存读取
        if cached_text is None:
            stats['cache_missing_skip'] += 1
            return None, None, None, False
        raw_text = cached_text
        stats['cache_used'] += 1
    else:
        # 缓存未命中，解析文件
        raw_text = extract_text_from_file(filepath_str)
        stats['cache_misses'] += 1
        
        # 立即保存到缓存（即使内容为空也保存，避免重复解析失败）
        if raw_text:
            cache.set(filepath_str, raw_text)
            stats['cache_saved'] += 1
        else:
            # 对于空内容，也保存空字符串，避免重复尝试解析
            cache.set(filepath_str, "", force_save=True)
            stats['cache_saved_empty'] += 1
    
    # 注意：raw_text 已经从缓存中获取，且已经在缓存阶段去除了人名
    
    if raw_text:
        cleaned = clean_text(raw_text)
        if cleaned:
            segmented = cut_words(cleaned)
            if segmented:
                stats['text_success'] += 1
                text_content = segmented
                if len(cleaned) < 20:
                    stats['text_short'] += 1
            else:
                text_content = "空文本"
        else:
            text_content = "空文本"
    else:
        text_content = "空文本"
        stats['text_empty'] += 1
    
    return text_content, filename_processed, original_filename, True


def extract_filename_features(filepath):
    """从文件路径中提取文件名特征"""
    filepath_obj = Path(filepath) if isinstance(filepath, str) else filepath
    filename = filepath_obj.stem
    
    cleaned = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', ' ', filename)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    
    return cleaned, filename


def load_data_with_cache(data_root):
    """带缓存的加载函数（支持PPTX和DOCX，立即保存缓存）"""
    texts = []
    filenames = []
    labels = []
    
    stats = {
        'total_files': 0,
        'extract_failed': 0,
        'text_empty': 0,
        'text_short': 0,
        'text_success': 0,
        'filename_success': 0,
        'augmented_count': 0,
        'success': 0,
        'cache_hits': 0,
        'cache_misses': 0,
        'cache_used': 0,
        'cache_saved': 0,
        'cache_saved_empty': 0,
        'cache_missing_skip': 0,
        'skipped_categories': 0,
        'by_category': {},
        'by_file_type': {}
    }
    
    cache = get_cache()
    skip_categories_set = set(SKIP_CATEGORIES)
    
    if skip_categories_set:
        print(f"\n🎯 跳过解析的科目（仅使用缓存）: {', '.join(skip_categories_set)}")
        print(f"   ⚠️  注意：这些科目的文件如果不在缓存中将被跳过")
    
    if REMOVE_PERSON_NAMES:
        print(f"\n👤 人名去除: 已启用（在缓存阶段使用jieba词性标注去除）")
    else:
        print(f"\n👤 人名去除: 已禁用")
    
    if not DOCX_SUPPORT:
        print(f"\n⚠️  警告: python-docx未安装，DOCX文件支持将受限")
        print(f"   请运行: pip install python-docx")
    
    print(f"\n📄 支持的文件格式: {', '.join(SUPPORTED_EXTENSIONS)}")
    
    for label_idx, cat in enumerate(CATEGORIES):
        cat_path = os.path.join(data_root, cat)
        if not os.path.isdir(cat_path):
            print(f"警告: 目录不存在 {cat_path}")
            stats['by_category'][cat] = {'total': 0, 'success': 0, 'augmented': 0, 'skipped': 0, 'cache_missing': 0}
            continue
        
        # 查找所有支持的文件类型
        files = []
        for ext in SUPPORTED_EXTENSIONS:
            files.extend(Path(cat_path).glob(f"*{ext}"))
        
        skip_parsing = cat in skip_categories_set
        if skip_parsing:
            stats['skipped_categories'] += 1
        
        # 按文件类型统计
        type_count = {}
        for f in files:
            ext = f.suffix.lower()
            type_count[ext] = type_count.get(ext, 0) + 1
        
        stats['by_category'][cat] = {
            'total': len(files), 
            'success': 0, 
            'augmented': 0,
            'skipped': 0 if not skip_parsing else len(files),
            'cache_missing': 0,
            'by_type': type_count
        }
        stats['total_files'] += len(files)
        
        if not files:
            print(f"警告: {cat} 目录下没有找到支持的文件")
            continue
        
        skip_msg = " [仅缓存模式]" if skip_parsing else ""
        print(f"\n处理 {cat} 目录{skip_msg}，共 {len(files)} 个文件")
        
        # 显示文件类型分布
        if type_count:
            type_str = ", ".join([f"{k}:{v}" for k, v in type_count.items()])
            print(f"  文件类型: {type_str}")
        
        for idx, filepath in enumerate(files):
            # 更新文件类型统计
            ext = filepath.suffix.lower()
            stats['by_file_type'][ext] = stats['by_file_type'].get(ext, 0) + 1
            
            text_content, filename_processed, original_filename, success = process_file(
                filepath, cat, label_idx, stats, skip_parsing=skip_parsing
            )
            
            if not success:
                stats['by_category'][cat]['cache_missing'] += 1
                if (idx + 1) % 20 == 0:
                    print(f"  ... 已处理 {idx + 1}/{len(files)} 个文件")
                continue
            
            filename_variants = [filename_processed]
            if ENABLE_FILENAME_AUGMENTATION and filename_processed:
                augmented_variants = augment_filename(original_filename)
                for variant in augmented_variants:
                    if variant and variant != original_filename:
                        variant_processed = process_filename_text(variant)
                        if variant_processed and variant_processed != filename_processed:
                            filename_variants.append(variant_processed)
                
                filename_variants = list(dict.fromkeys(filename_variants))[:FILENAME_AUGMENTATION_COUNT]
            
            for fname_var in filename_variants:
                if fname_var:
                    texts.append(text_content)
                    filenames.append(fname_var)
                    labels.append(label_idx)
                    stats['success'] += 1
                    stats['by_category'][cat]['success'] += 1
                    
                    if fname_var != filename_processed:
                        stats['augmented_count'] += 1
                        stats['by_category'][cat]['augmented'] += 1
            
            if (idx + 1) % 20 == 0:
                print(f"  ... 已处理 {idx + 1}/{len(files)} 个文件")
                # 每处理20个文件，显示一次缓存状态
                cache_stats = cache.get_stats()
                print(f"     缓存: {cache_stats['cached_files']} 个文件, {cache_stats['non_empty']} 非空")
        
        cat_stats = stats['by_category'][cat]
        if skip_parsing:
            print(f"  {cat}: 原始{cat_stats['total']}个文件 → "
                  f"成功{cat_stats['success']}个样本（含{cat_stats['augmented']}个增强），"
                  f"缓存缺失{cat_stats['cache_missing']}个文件")
        else:
            print(f"  {cat}: 原始{cat_stats['total']}个文件 → {cat_stats['success']}个样本（含{cat_stats['augmented']}个增强）")
    
    return texts, filenames, labels, stats


def print_stats(stats):
    """打印详细统计信息"""
    print("\n" + "="*50)
    print("数据加载统计报告（带缓存）")
    print("="*50)
    print(f"总文件数:        {stats['total_files']}")
    print(f"生成样本数:      {stats['success']}")
    print(f"增强样本数:      {stats['augmented_count']}")
    if stats['success'] > 0:
        print(f"  └─ 增强比例:   {stats['augmented_count']/stats['success']*100:.1f}%")
    
    print(f"\n缓存统计:")
    print(f"  缓存命中:      {stats.get('cache_hits', 0)}")
    print(f"  缓存未命中:    {stats.get('cache_misses', 0)}")
    print(f"  缓存使用:      {stats.get('cache_used', 0)}")
    print(f"  缓存保存:      {stats.get('cache_saved', 0)}")
    print(f"  缓存保存空:    {stats.get('cache_saved_empty', 0)}")
    
    if stats.get('cache_missing_skip', 0) > 0:
        print(f"\n⚠️  跳过解析模式中因缓存缺失跳过的文件: {stats['cache_missing_skip']}")
    
    print(f"\n文本状态:")
    print(f"  有文本内容:    {stats['text_success']}")
    print(f"  文本过短:      {stats['text_short']}")
    print(f"  文本为空:      {stats['text_empty']}")
    
    # 文件类型分布
    if stats.get('by_file_type'):
        print(f"\n文件类型分布:")
        for ftype, count in sorted(stats['by_file_type'].items()):
            print(f"  {ftype}: {count} 个文件")
    
    cache = get_cache()
    cache_stats = cache.get_stats()
    print(f"\n缓存详细统计:")
    print(f"  总缓存条目:    {cache_stats['cached_files']}")
    print(f"  非空条目:      {cache_stats['non_empty']}")
    print(f"  空条目:        {cache_stats['empty_entries']}")
    print(f"  缓存文件大小:  {cache_stats['cache_file_size'] / 1024:.1f} KB")
    
    # 打印人名去除统计
    if REMOVE_PERSON_NAMES:
        print(f"\n👤 人名去除统计:")
        print(f"  总共去除字符:  {cache_stats['names_removed_total']}")
        print(f"  涉及文件数:    {cache_stats['files_with_names_removed']}")
    
    print("\n各类别统计:")
    for cat, info in stats['by_category'].items():
        if info['total'] > 0:
            # 显示文件类型分布
            type_info = ""
            if info.get('by_type'):
                type_info = f" [{', '.join([f'{k}:{v}' for k, v in info['by_type'].items()])}]"
            if info.get('cache_missing', 0) > 0:
                print(f"  {cat}{type_info}: {info['total']}个文件 → {info['success']}个样本（{info['augmented']}增强，{info['cache_missing']}缓存缺失）")
            else:
                print(f"  {cat}{type_info}: {info['total']}个文件 → {info['success']}个样本（{info['augmented']}增强）")
    print("="*50)


def build_optimized_multimodal_model(vocab_size, max_seq_len, max_filename_len, 
                                      filename_vocab_size, num_classes):
    """
    构建优化版的多模态模型
    特点：
    1. 增加kernel_size=2的卷积核（捕获更短模式）
    2. 降低融合层Dropout率（保留更多信息）
    3. 增加文本分支复杂度
    """
    
    # 文本分支 - 增强版
    text_input = Input(shape=(max_seq_len,), name='text_input')
    text_embedding = Embedding(vocab_size, EMBEDDING_DIM, name='text_embedding')(text_input)
    
    # 多尺度卷积（增加kernel_size=2）
    conv_2 = Conv1D(filters=128, kernel_size=2, activation='relu', padding='same', name='text_conv_2')(text_embedding)
    conv_3 = Conv1D(filters=128, kernel_size=3, activation='relu', padding='same', name='text_conv_3')(text_embedding)
    conv_4 = Conv1D(filters=128, kernel_size=4, activation='relu', padding='same', name='text_conv_4')(text_embedding)
    conv_5 = Conv1D(filters=128, kernel_size=5, activation='relu', padding='same', name='text_conv_5')(text_embedding)
    
    pool_2 = GlobalMaxPooling1D(name='text_pool_2')(conv_2)
    pool_3 = GlobalMaxPooling1D(name='text_pool_3')(conv_3)
    pool_4 = GlobalMaxPooling1D(name='text_pool_4')(conv_4)
    pool_5 = GlobalMaxPooling1D(name='text_pool_5')(conv_5)
    
    # 拼接所有尺度特征
    text_multi = Concatenate(name='text_multi')([pool_2, pool_3, pool_4, pool_5])
    text_dense = Dense(128, activation='relu', name='text_dense')(text_multi)
    text_dropout = Dropout(0.5, name='text_dropout')(text_dense)
    
    # 文件名分支
    filename_input = Input(shape=(max_filename_len,), name='filename_input')
    filename_embedding = Embedding(filename_vocab_size, FILENAME_EMBEDDING_DIM, name='filename_embedding')(filename_input)
    filename_conv = Conv1D(filters=64, kernel_size=2, activation='relu', name='filename_conv')(filename_embedding)
    filename_pool = GlobalMaxPooling1D(name='filename_pool')(filename_conv)
    filename_dense = Dense(32, activation='relu', name='filename_dense')(filename_pool)
    filename_dropout = Dropout(0.3, name='filename_dropout')(filename_dense)
    
    # 融合层 - 降低Dropout率
    combined = Concatenate(name='combined')([text_dropout, filename_dropout])
    final_dense = Dense(64, activation='relu', name='final_dense')(combined)
    final_dropout = Dropout(0.3, name='final_dropout')(final_dense)  # 0.5 → 0.3
    output = Dense(num_classes, activation='softmax', name='output')(final_dropout)
    
    model = Model(inputs=[text_input, filename_input], outputs=output)
    model.compile(loss='categorical_crossentropy', optimizer='adam', metrics=['accuracy'])
    
    return model


def save_tokenizers_without_keras(text_tokenizer, filename_tokenizer, categories, config):
    """
    保存无Keras依赖的Tokenizer文件（用于推理）
    在20000词处截断（只保留前MAX_NB_WORDS个词）
    """
    print("\n" + "="*50)
    print("保存无依赖的Tokenizer文件（词表截断）")
    print("="*50)
    
    # 处理文本Tokenizer - 截断到MAX_NB_WORDS
    full_word_index = text_tokenizer.word_index
    full_vocab_size = len(full_word_index)
    
    # 截断：只保留前MAX_NB_WORDS个最常见的词
    truncated_word_index = {}
    truncated_index_word = {}
    
    for i, (word, idx) in enumerate(full_word_index.items()):
        if i < MAX_NB_WORDS:
            truncated_word_index[word] = idx
            truncated_index_word[idx] = word
    
    actual_vocab_size = len(truncated_word_index)
    
    print(f"  文本Tokenizer:")
    print(f"    - 原始词表大小: {full_vocab_size}")
    print(f"    - 截断后词表大小: {actual_vocab_size} (限制在 {MAX_NB_WORDS})")
    print(f"    - 截断比例: {(1 - actual_vocab_size/full_vocab_size)*100:.1f}%")
    
    text_tokenizer_data = {
        'word_index': truncated_word_index,
        'index_word': truncated_index_word,
        'word_counts': {word: text_tokenizer.word_counts.get(word, 0) 
                       for word in truncated_word_index.keys()},
        'document_count': text_tokenizer.document_count,
        'vocab_size': actual_vocab_size,
        'max_words': MAX_NB_WORDS,
        'is_truncated': True,
        'original_vocab_size': full_vocab_size
    }
    
    # 处理文件名Tokenizer - 截断到MAX_FILENAME_WORDS
    full_filename_word_index = filename_tokenizer.word_index
    full_filename_vocab_size = len(full_filename_word_index)
    
    truncated_filename_word_index = {}
    truncated_filename_index_word = {}
    
    for i, (word, idx) in enumerate(full_filename_word_index.items()):
        if i < MAX_FILENAME_WORDS:
            truncated_filename_word_index[word] = idx
            truncated_filename_index_word[idx] = word
    
    actual_filename_vocab_size = len(truncated_filename_word_index)
    
    print(f"\n  文件名Tokenizer:")
    print(f"    - 原始词表大小: {full_filename_vocab_size}")
    print(f"    - 截断后词表大小: {actual_filename_vocab_size} (限制在 {MAX_FILENAME_WORDS})")
    print(f"    - 截断比例: {(1 - actual_filename_vocab_size/full_filename_vocab_size)*100:.1f}%")
    
    filename_tokenizer_data = {
        'word_index': truncated_filename_word_index,
        'index_word': truncated_filename_index_word,
        'word_counts': {word: filename_tokenizer.word_counts.get(word, 0)
                       for word in truncated_filename_word_index.keys()},
        'document_count': filename_tokenizer.document_count,
        'vocab_size': actual_filename_vocab_size,
        'max_words': MAX_FILENAME_WORDS,
        'is_truncated': True,
        'original_vocab_size': full_filename_vocab_size
    }
    
    # 保存文件
    with open('text_tokenizer_none.pkl', 'wb') as f:
        pickle.dump(text_tokenizer_data, f)
    print(f"\n✅ 已保存: text_tokenizer_none.pkl")
    
    with open('filename_tokenizer_none.pkl', 'wb') as f:
        pickle.dump(filename_tokenizer_data, f)
    print(f"✅ 已保存: filename_tokenizer_none.pkl")
    
    # 保存类别和配置
    with open('categories.pkl', 'wb') as f:
        pickle.dump(categories, f)
    
    # 更新配置，记录截断信息
    config['max_nb_words'] = MAX_NB_WORDS
    config['max_filename_words'] = MAX_FILENAME_WORDS
    config['actual_text_vocab_size'] = actual_vocab_size
    config['actual_filename_vocab_size'] = actual_filename_vocab_size
    
    with open('config_optimized.pkl', 'wb') as f:
        pickle.dump(config, f)
    
    print(f"✅ 已保存: categories.pkl")
    print(f"✅ 已保存: config_optimized.pkl")
    
    # 打印汇总信息
    print("\n" + "="*50)
    print("📊 词表截断汇总:")
    print(f"  文本: {full_vocab_size} → {actual_vocab_size} (保留前{MAX_NB_WORDS}词)")
    print(f"  文件名: {full_filename_vocab_size} → {actual_filename_vocab_size} (保留前{MAX_FILENAME_WORDS}词)")
    print(f"  建议 embedding_dim: {min(EMBEDDING_DIM, actual_vocab_size, actual_filename_vocab_size)}")
    print("="*50)
    
    return text_tokenizer_data, filename_tokenizer_data


def save_json_vocabularies(text_tokenizer, filename_tokenizer):
    """
    保存JSON格式的词表文件（截断在MAX_NB_WORDS和MAX_FILENAME_WORDS）
    便于外部工具读取和使用
    """
    print("\n" + "="*50)
    print("保存JSON格式词表文件")
    print("="*50)
    
    # ========== 1. 保存文本词表 ==========
    full_word_index = text_tokenizer.word_index
    
    # 截断：只保留前MAX_NB_WORDS个最常见的词
    truncated_word_index = {}
    truncated_word_counts = {}
    
    for i, (word, idx) in enumerate(full_word_index.items()):
        if i < MAX_NB_WORDS:
            truncated_word_index[word] = idx
            truncated_word_counts[word] = text_tokenizer.word_counts.get(word, 0)
    
    # 构建文本词表JSON
    text_vocab_json = {
        "metadata": {
            "type": "text_vocabulary",
            "version": CACHE_VERSION,
            "max_words": MAX_NB_WORDS,
            "original_size": len(full_word_index),
            "actual_size": len(truncated_word_index),
            "is_truncated": True,
            "created_at": datetime.now().isoformat(),
            "description": "文本分类器的词表，按词频降序排列，仅保留前{}个最常见的词".format(MAX_NB_WORDS)
        },
        "vocabulary": [
            {
                "index": idx,
                "word": word,
                "frequency": truncated_word_counts.get(word, 0)
            }
            for word, idx in truncated_word_index.items()
        ],
        "word_to_index": truncated_word_index,
        "index_to_word": {str(idx): word for word, idx in truncated_word_index.items()}
    }
    
    text_vocab_json["vocabulary"].sort(key=lambda x: x["index"])
    
    with open('text_vocabulary.json', 'w', encoding='utf-8') as f:
        json.dump(text_vocab_json, f, ensure_ascii=False, indent=2)
    print(f"✅ 已保存: text_vocabulary.json")
    print(f"   - 词表大小: {len(truncated_word_index)} 词")
    print(f"   - 原始大小: {len(full_word_index)} 词")
    print(f"   - 文件大小: {os.path.getsize('text_vocabulary.json') / 1024:.1f} KB")
    
    # 简洁版本
    simple_text_vocab = {
        "metadata": {
            "type": "text_vocabulary_simple",
            "max_words": MAX_NB_WORDS,
            "actual_size": len(truncated_word_index)
        },
        "words_by_frequency": list(truncated_word_index.keys()),
        "frequency_map": truncated_word_counts
    }
    
    with open('text_vocabulary_simple.json', 'w', encoding='utf-8') as f:
        json.dump(simple_text_vocab, f, ensure_ascii=False, indent=2)
    print(f"✅ 已保存: text_vocabulary_simple.json")
    
    # ========== 2. 保存文件名词表 ==========
    full_filename_word_index = filename_tokenizer.word_index
    
    truncated_filename_word_index = {}
    truncated_filename_word_counts = {}
    
    for i, (word, idx) in enumerate(full_filename_word_index.items()):
        if i < MAX_FILENAME_WORDS:
            truncated_filename_word_index[word] = idx
            truncated_filename_word_counts[word] = filename_tokenizer.word_counts.get(word, 0)
    
    filename_vocab_json = {
        "metadata": {
            "type": "filename_vocabulary",
            "version": CACHE_VERSION,
            "max_words": MAX_FILENAME_WORDS,
            "original_size": len(full_filename_word_index),
            "actual_size": len(truncated_filename_word_index),
            "is_truncated": True,
            "created_at": datetime.now().isoformat(),
            "description": "文件名词表，按词频降序排列，仅保留前{}个最常见的词".format(MAX_FILENAME_WORDS)
        },
        "vocabulary": [
            {
                "index": idx,
                "word": word,
                "frequency": truncated_filename_word_counts.get(word, 0)
            }
            for word, idx in truncated_filename_word_index.items()
        ],
        "word_to_index": truncated_filename_word_index,
        "index_to_word": {str(idx): word for word, idx in truncated_filename_word_index.items()}
    }
    
    filename_vocab_json["vocabulary"].sort(key=lambda x: x["index"])
    
    with open('filename_vocabulary.json', 'w', encoding='utf-8') as f:
        json.dump(filename_vocab_json, f, ensure_ascii=False, indent=2)
    print(f"\n✅ 已保存: filename_vocabulary.json")
    print(f"   - 词表大小: {len(truncated_filename_word_index)} 词")
    print(f"   - 原始大小: {len(full_filename_word_index)} 词")
    print(f"   - 文件大小: {os.path.getsize('filename_vocabulary.json') / 1024:.1f} KB")
    
    # 简洁版本
    simple_filename_vocab = {
        "metadata": {
            "type": "filename_vocabulary_simple",
            "max_words": MAX_FILENAME_WORDS,
            "actual_size": len(truncated_filename_word_index)
        },
        "words_by_frequency": list(truncated_filename_word_index.keys()),
        "frequency_map": truncated_filename_word_counts
    }
    
    with open('filename_vocabulary_simple.json', 'w', encoding='utf-8') as f:
        json.dump(simple_filename_vocab, f, ensure_ascii=False, indent=2)
    print(f"✅ 已保存: filename_vocabulary_simple.json")
    
    # ========== 3. 保存词频统计报告 ==========
    word_frequencies = [(word, count) for word, count in truncated_word_counts.items()]
    word_frequencies.sort(key=lambda x: x[1], reverse=True)
    
    top_100_words = [
        {"rank": i+1, "word": word, "frequency": freq}
        for i, (word, freq) in enumerate(word_frequencies[:100])
    ]
    
    freq_buckets = {
        "1-10": 0,
        "11-50": 0,
        "51-100": 0,
        "101-500": 0,
        "501-1000": 0,
        "1001+": 0
    }
    
    for _, freq in word_frequencies:
        if freq <= 10:
            freq_buckets["1-10"] += 1
        elif freq <= 50:
            freq_buckets["11-50"] += 1
        elif freq <= 100:
            freq_buckets["51-100"] += 1
        elif freq <= 500:
            freq_buckets["101-500"] += 1
        elif freq <= 1000:
            freq_buckets["501-1000"] += 1
        else:
            freq_buckets["1001+"] += 1
    
    word_frequency_report = {
        "metadata": {
            "type": "word_frequency_report",
            "created_at": datetime.now().isoformat(),
            "vocabulary_size": len(truncated_word_index),
            "total_word_count": sum(truncated_word_counts.values())
        },
        "top_100_words": top_100_words,
        "frequency_distribution": freq_buckets,
        "statistics": {
            "min_frequency": min(truncated_word_counts.values()) if truncated_word_counts else 0,
            "max_frequency": max(truncated_word_counts.values()) if truncated_word_counts else 0,
            "avg_frequency": sum(truncated_word_counts.values()) / len(truncated_word_counts) if truncated_word_counts else 0,
            "median_frequency": sorted(truncated_word_counts.values())[len(truncated_word_counts)//2] if truncated_word_counts else 0
        }
    }
    
    with open('word_frequency_report.json', 'w', encoding='utf-8') as f:
        json.dump(word_frequency_report, f, ensure_ascii=False, indent=2)
    print(f"\n✅ 已保存: word_frequency_report.json")
    
    # ========== 4. 保存类别映射 ==========
    category_mapping = {
        "metadata": {
            "type": "category_mapping",
            "created_at": datetime.now().isoformat(),
            "num_categories": len(CATEGORIES)
        },
        "categories": CATEGORIES,
        "index_to_category": {str(i): cat for i, cat in enumerate(CATEGORIES)},
        "category_to_index": {cat: i for i, cat in enumerate(CATEGORIES)}
    }
    
    with open('category_mapping.json', 'w', encoding='utf-8') as f:
        json.dump(category_mapping, f, ensure_ascii=False, indent=2)
    print(f"✅ 已保存: category_mapping.json")
    
    print("\n" + "="*50)
    print("📊 JSON词表文件汇总:")
    print(f"  文本词表 (完整): text_vocabulary.json ({len(truncated_word_index)} 词)")
    print(f"  文本词表 (简洁): text_vocabulary_simple.json")
    print(f"  文件名词表 (完整): filename_vocabulary.json ({len(truncated_filename_word_index)} 词)")
    print(f"  文件名词表 (简洁): filename_vocabulary_simple.json")
    print(f"  词频统计报告: word_frequency_report.json")
    print(f"  类别映射: category_mapping.json")
    print("="*50)
    
    return True


# ---------- 主流程 ----------
def main():
    print("="*50)
    print("文件学科分类器训练（优化版 - 缓存修复）")
    print(f"支持类别: {', '.join(CATEGORIES)}")
    print(f"支持格式: {', '.join(SUPPORTED_EXTENSIONS)}")
    print(f"文本权重: {TEXT_WEIGHT}, 文件名权重: {FILENAME_WEIGHT}")
    print(f"数据增强: {'启用' if ENABLE_FILENAME_AUGMENTATION else '禁用'}")
    print(f"缓存: {'启用' if ENABLE_CACHE else '禁用'} (版本: {CACHE_VERSION})")
    print(f"类别权重: {'启用' if USE_CLASS_WEIGHTS else '禁用'}")
    print(f"人名去除: {'启用 (缓存阶段)' if REMOVE_PERSON_NAMES else '禁用'}")
    print(f"词表截断: 限制在 {MAX_NB_WORDS} 词")
    print(f"DOCX支持: {'✅ 已启用' if DOCX_SUPPORT else '❌ 未安装'}")
    
    if SKIP_CATEGORIES:
        print(f"🎯 仅缓存模式科目: {', '.join(SKIP_CATEGORIES)}")
    
    print("="*50)
    
    if not os.path.exists(DATA_ROOT):
        print(f"错误: 数据目录不存在 '{DATA_ROOT}'")
        return
    
    if FORCE_REFRESH_CACHE:
        print("\n强制刷新缓存模式...")
        cache = get_cache()
        cache.clear()
    
    # 1. 加载数据
    print("\n步骤1: 加载文件（使用缓存加速，立即保存）...")
    import time
    start_time = time.time()
    texts, filenames, labels, stats = load_data_with_cache(DATA_ROOT)
    load_time = time.time() - start_time
    print(f"数据加载耗时: {load_time:.2f} 秒")
    
    print_stats(stats)
    
    # 打印人名去除详情
    cache = get_cache()
    cache.print_name_removal_summary()
    cache.print_cache_summary()
    
    # 最终保存缓存（确保所有数据都已保存）
    print("\n💾 最终保存缓存...")
    cache.save()
    
    if len(texts) == 0:
        print("\n错误: 未找到任何有效文件！")
        print("请确保数据目录结构为: data/类别/*.pptx 或 data/类别/*.docx")
        return
    
    # 2. 划分数据集
    print("\n步骤2: 划分数据集...")
    X_train_text, X_test_text, X_train_filename, X_test_filename, y_train, y_test = train_test_split(
        texts, filenames, labels, test_size=TEST_SPLIT, stratify=labels, random_state=42
    )
    
    val_ratio = VALIDATION_SPLIT / (1 - TEST_SPLIT)
    X_train_text, X_val_text, X_train_filename, X_val_filename, y_train, y_val = train_test_split(
        X_train_text, X_train_filename, y_train, test_size=val_ratio, 
        stratify=y_train, random_state=42
    )
    print(f"训练集: {len(X_train_text)}")
    print(f"验证集: {len(X_val_text)}")
    print(f"测试集: {len(X_test_text)}")
    
    # 3. 构建Tokenizer（使用截断）
    print("\n步骤3: 构建词表并转换序列...")
    print(f"  词表截断: 最多保留 {MAX_NB_WORDS} 个最常见的词")
    
    text_tokenizer = Tokenizer(num_words=MAX_NB_WORDS, oov_token='<UNK>')
    text_tokenizer.fit_on_texts(X_train_text)
    
    text_vocab_size = min(MAX_NB_WORDS, len(text_tokenizer.word_index) + 1)
    
    print(f"  文本原始词汇量: {len(text_tokenizer.word_index)}")
    print(f"  文本实际使用词汇量: {text_vocab_size} (截断在 {MAX_NB_WORDS})")
    
    filename_tokenizer = Tokenizer(num_words=MAX_FILENAME_WORDS, oov_token='<UNK>')
    filename_tokenizer.fit_on_texts(X_train_filename)
    
    filename_vocab_size = min(MAX_FILENAME_WORDS, len(filename_tokenizer.word_index) + 1)
    
    print(f"  文件名字典大小: {len(filename_tokenizer.word_index)}")
    print(f"  文件名实际使用词汇量: {filename_vocab_size} (截断在 {MAX_FILENAME_WORDS})")
    
    def seq_and_pad_text(texts):
        seqs = text_tokenizer.texts_to_sequences(texts)
        return pad_sequences(seqs, maxlen=MAX_SEQUENCE_LENGTH, padding='post', truncating='post')
    
    def seq_and_pad_filename(filenames):
        seqs = filename_tokenizer.texts_to_sequences(filenames)
        return pad_sequences(seqs, maxlen=MAX_FILENAME_LENGTH, padding='post', truncating='post')
    
    X_train_text_pad = seq_and_pad_text(X_train_text)
    X_val_text_pad = seq_and_pad_text(X_val_text)
    X_test_text_pad = seq_and_pad_text(X_test_text)
    
    X_train_filename_pad = seq_and_pad_filename(X_train_filename)
    X_val_filename_pad = seq_and_pad_filename(X_val_filename)
    X_test_filename_pad = seq_and_pad_filename(X_test_filename)
    
    y_train_cat = to_categorical(y_train, num_classes=len(CATEGORIES))
    y_val_cat = to_categorical(y_val, num_classes=len(CATEGORIES))
    y_test_cat = to_categorical(y_test, num_classes=len(CATEGORIES))
    
    print(f"文本输入形状: {X_train_text_pad.shape}")
    print(f"文件名输入形状: {X_train_filename_pad.shape}")
    
    # 4. 构建模型
    print("\n步骤4: 构建优化版多模态模型...")
    model = build_optimized_multimodal_model(
        vocab_size=text_vocab_size,
        max_seq_len=MAX_SEQUENCE_LENGTH,
        max_filename_len=MAX_FILENAME_LENGTH,
        filename_vocab_size=filename_vocab_size,
        num_classes=len(CATEGORIES)
    )
    model.summary()
    
    # 5. 计算类别权重（可选）
    class_weights = None
    if USE_CLASS_WEIGHTS:
        print("\n计算类别权重...")
        class_weights = compute_class_weight(
            'balanced',
            classes=np.unique(y_train),
            y=y_train
        )
        class_weights = dict(enumerate(class_weights))
        for i, cat in enumerate(CATEGORIES):
            print(f"  {cat}: {class_weights[i]:.4f}")
    
    # 6. 训练
    print("\n步骤5: 开始训练...")
    early_stop = EarlyStopping(
        monitor='val_loss', 
        patience=6,
        restore_best_weights=True,
        min_delta=0.0001
    )
    checkpoint = ModelCheckpoint(
        'best_model_optimized.keras', 
        monitor='val_accuracy', 
        save_best_only=True, 
        mode='max'
    )
    reduce_lr = ReduceLROnPlateau(
        monitor='val_loss', 
        factor=0.5, 
        patience=3, 
        min_lr=1e-6,
        cooldown=1
    )
    
    history = model.fit(
        [X_train_text_pad, X_train_filename_pad], y_train_cat,
        batch_size=BATCH_SIZE,
        epochs=EPOCHS,
        validation_data=([X_val_text_pad, X_val_filename_pad], y_val_cat),
        callbacks=[early_stop, checkpoint, reduce_lr],
        class_weight=class_weights,
        verbose=1
    )
    
    # 7. 评估
    print("\n步骤6: 模型评估...")
    loss, acc = model.evaluate([X_test_text_pad, X_test_filename_pad], y_test_cat, verbose=0)
    print(f"测试集准确率: {acc:.4f}")
    
    y_pred_proba = model.predict([X_test_text_pad, X_test_filename_pad])
    y_pred = np.argmax(y_pred_proba, axis=1)
    print("\n详细分类报告:")
    print(classification_report(y_test, y_pred, target_names=CATEGORIES))
    print("混淆矩阵:")
    print(confusion_matrix(y_test, y_pred))
    
    # 8. 保存模型和Tokenizer
    print("\n步骤7: 保存模型...")
    model.save('textcnn_optimized_classifier.keras')
    
    with open('text_tokenizer.pkl', 'wb') as f:
        pickle.dump(text_tokenizer, f)
    with open('filename_tokenizer.pkl', 'wb') as f:
        pickle.dump(filename_tokenizer, f)
    
    config = {
        'max_sequence_length': MAX_SEQUENCE_LENGTH,
        'max_filename_length': MAX_FILENAME_LENGTH,
        'categories': CATEGORIES,
        'text_vocab_size': text_vocab_size,
        'filename_vocab_size': filename_vocab_size,
        'embedding_dim': EMBEDDING_DIM,
        'filename_embedding_dim': FILENAME_EMBEDDING_DIM,
        'text_weight': TEXT_WEIGHT,
        'filename_weight': FILENAME_WEIGHT,
        'cache_version': CACHE_VERSION,
        'skip_categories': SKIP_CATEGORIES,
        'use_class_weights': USE_CLASS_WEIGHTS,
        'remove_person_names': REMOVE_PERSON_NAMES,
        'max_nb_words': MAX_NB_WORDS,
        'max_filename_words': MAX_FILENAME_WORDS,
        'supported_extensions': SUPPORTED_EXTENSIONS
    }
    
    save_tokenizers_without_keras(text_tokenizer, filename_tokenizer, CATEGORIES, config)
    save_json_vocabularies(text_tokenizer, filename_tokenizer)
    
    print("\n" + "="*50)
    print("训练完成！已保存以下文件:")
    print("  - textcnn_optimized_classifier.keras")
    print("  - best_model_optimized.keras")
    print("  - text_tokenizer.pkl (完整版)")
    print("  - filename_tokenizer.pkl (完整版)")
    print("  - text_tokenizer_none.pkl (无依赖版，截断)")
    print("  - filename_tokenizer_none.pkl (无依赖版，截断)")
    print("  - categories.pkl")
    print("  - config_optimized.pkl")
    print("\n📝 JSON词表文件:")
    print("  - text_vocabulary.json (完整文本词表)")
    print("  - text_vocabulary_simple.json (简洁文本词表)")
    print("  - filename_vocabulary.json (完整文件名词表)")
    print("  - filename_vocabulary_simple.json (简洁文件名词表)")
    print("  - word_frequency_report.json (词频统计报告)")
    print("  - category_mapping.json (类别映射)")
    print(f"\n缓存目录: {CACHE_DIR}/")
    print(f"  - file_cache_{CACHE_VERSION}.json")
    print(f"  - cache_metadata_{CACHE_VERSION}.json")
    print(f"  - cache_debug_{CACHE_VERSION}.log")
    if REMOVE_PERSON_NAMES:
        print(f"  - name_removal_stats_{CACHE_VERSION}.json")
    print("="*50)
    
    # 可选：保存训练历史
    try:
        import matplotlib.pyplot as plt
        fig, axes = plt.subplots(1, 2, figsize=(12, 4))
        
        axes[0].plot(history.history['loss'], label='train')
        axes[0].plot(history.history['val_loss'], label='val')
        axes[0].set_title('Model Loss')
        axes[0].set_xlabel('Epoch')
        axes[0].set_ylabel('Loss')
        axes[0].legend()
        
        axes[1].plot(history.history['accuracy'], label='train')
        axes[1].plot(history.history['val_accuracy'], label='val')
        axes[1].set_title('Model Accuracy')
        axes[1].set_xlabel('Epoch')
        axes[1].set_ylabel('Accuracy')
        axes[1].legend()
        
        plt.tight_layout()
        plt.savefig('training_history.png')
        print("训练曲线已保存: training_history.png")
        plt.close()
    except:
        pass


if __name__ == "__main__":
    random.seed(42)
    np.random.seed(42)
    main()