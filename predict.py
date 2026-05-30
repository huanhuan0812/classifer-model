#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX学科分类器预测脚本 - 内存优化版
支持7类别：语文、数学、英语、物理、化学、生物、班会
解决了交互式模式下内存持续增长的问题
"""

import os
import sys
import gc
import warnings
import re
import pickle
import argparse
import logging
from pathlib import Path
from datetime import datetime

# 屏蔽所有警告
warnings.filterwarnings('ignore')
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'  # 屏蔽TensorFlow日志
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'  # 屏蔽oneDNN警告

# 设置logging级别
logging.getLogger().setLevel(logging.ERROR)

# 导入tensorflow并配置内存优化
import tensorflow as tf

# 屏蔽TensorFlow的GPU警告
os.environ['TF_GPU_ALLOCATOR'] = 'cuda_malloc_async'

# 设置TensorFlow日志级别
tf.get_logger().setLevel('ERROR')
tf.autograph.set_verbosity(0)

# 配置CPU线程数（减少内存碎片）
try:
    tf.config.threading.set_intra_op_parallelism_threads(2)
    tf.config.threading.set_inter_op_parallelism_threads(2)
except:
    pass

# 设置内存增长
try:
    gpus = tf.config.experimental.list_physical_devices('GPU')
    if gpus:
        for gpu in gpus:
            tf.config.experimental.set_memory_growth(gpu, True)
except:
    pass

from tensorflow.keras.models import load_model
from tensorflow.keras.preprocessing.sequence import pad_sequences
import jieba
from pptx import Presentation

# ---------- 配置 ----------
MODEL_PATH = "textcnn_balanced_classifier.keras"
TEXT_TOKENIZER_PATH = "text_tokenizer.pkl"
FILENAME_TOKENIZER_PATH = "filename_tokenizer.pkl"
CATEGORIES_PATH = "categories.pkl"
CONFIG_PATH = "config_balanced.pkl"

# 停用词表（与训练时保持一致）
STOPWORDS = set([
    '的', '了', '是', '我', '你', '他', '她', '它', '我们', '你们', '他们',
    '这', '那', '有', '在', '不', '和', '与', '就', '都', '而', '及', '或',
    '一个', '这个', '那个', '那些', '这些', '这里', '那里', '然后', '因为',
    '所以', '但是', '如果', '虽然', '然而', '并且', '或者'
])

# 有意义的单字词
MEANINGFUL_SINGLE_CHARS = {'圆', '力', '氧', '氢', '碳', '钠', '酸', '碱', '盐', 
                           '电', '光', '声', '热', '诗', '词', '歌', '曲', '数',
                           '方', '程', '函', '数', '角', '形', '体', '积'}

# 抑制jieba输出
jieba.setLogLevel(logging.ERROR)

# ---------- 全局变量 ----------
_model = None
_text_tokenizer = None
_filename_tokenizer = None
_categories = None
_config = None

# 预测计数器
_prediction_count = 0


# ---------- 模型加载 ----------
def load_models():
    """加载模型和分词器"""
    global _model, _text_tokenizer, _filename_tokenizer, _categories, _config
    
    print("正在加载模型...", end=' ', flush=True)
    try:
        _model = load_model(MODEL_PATH, compile=False)
        # 修复：新版本TensorFlow不需要调用 _make_predict_function
        # 只需进行一次空预测来初始化模型
        print("✓", end=' ')
    except Exception as e:
        print(f"✗\n错误: {e}")
        raise
    
    print("加载文本分词器...", end=' ', flush=True)
    with open(TEXT_TOKENIZER_PATH, 'rb') as f:
        _text_tokenizer = pickle.load(f)
    print("✓", end=' ')
    
    print("加载文件名词典...", end=' ', flush=True)
    with open(FILENAME_TOKENIZER_PATH, 'rb') as f:
        _filename_tokenizer = pickle.load(f)
    print("✓", end=' ')
    
    print("加载类别映射...", end=' ', flush=True)
    with open(CATEGORIES_PATH, 'rb') as f:
        _categories = pickle.load(f)
    print("✓", end=' ')
    
    print("加载配置...", end=' ', flush=True)
    with open(CONFIG_PATH, 'rb') as f:
        _config = pickle.load(f)
    print("✓")
    
    print(f"\n模型已就绪！支持 {len(_categories)} 个类别: {', '.join(_categories)}\n")
    
    # 预热模型（进行一次空预测，避免第一次预测时的延迟）
    try:
        import numpy as np
        dummy_text = np.zeros((1, _config['max_sequence_length']), dtype=np.int32)
        dummy_filename = np.zeros((1, _config['max_filename_length']), dtype=np.int32)
        _model.predict([dummy_text, dummy_filename], verbose=0)
        print("模型预热完成\n")
    except:
        pass


# ---------- PPTX 文本提取（优化内存）----------
def extract_text_from_pptx(pptx_path):
    """从PPTX文件中提取所有文本内容（优化内存释放）"""
    text_parts = []
    prs = None
    
    try:
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text.strip())
                    
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text_parts.append(cell.text.strip())
                    
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text:
                                text_parts.append(paragraph.text.strip())
                except:
                    continue
            
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        text_parts.append(notes_slide.notes_text_frame.text.strip())
            except:
                pass
            
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    text_parts.append(slide.shapes.title.text.strip())
            except:
                pass
            
            # 主动释放slide对象
            del slide
        
        result = " ".join(text_parts)
        
    except Exception as e:
        result = ""
    finally:
        # 确保释放Presentation对象
        if prs is not None:
            del prs
        
        # 清空临时列表
        text_parts.clear()
        del text_parts
    
    return result


# ---------- 文本预处理 ----------
def clean_text(text):
    """基础清洗：保留中文、英文、数字、常用标点"""
    if not text:
        return ""
    
    text = re.sub(r'http\S+|www\S+|https\S+', '', text)
    text = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9，。！？；：""''（）【】《》、 ]', '', text)
    text = re.sub(r'\s+', ' ', text)
    
    return text.strip()


def cut_words(text):
    """使用jieba进行分词（不缓存结果，减少内存）"""
    if not text:
        return ""
    
    # 使用jieba的精确模式，HMM=False减少计算
    words = jieba.cut(text, HMM=False)
    filtered = []
    for w in words:
        if not w or w.strip() == '':
            continue
        if w not in STOPWORDS:
            filtered.append(w)
    
    if not filtered:
        # 如果全部被过滤，返回原始分词（不过滤停用词）
        words = jieba.cut(text, HMM=False)
        filtered = [w for w in words if w.strip()]
    
    if not filtered:
        return "空文本"
    
    result = " ".join(filtered)
    
    # 清空临时列表
    filtered.clear()
    del filtered
    
    return result


def extract_filename_features(filepath):
    """从文件路径中提取文件名特征"""
    filepath_obj = Path(filepath) if isinstance(filepath, str) else filepath
    filename = filepath_obj.stem
    
    cleaned = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', ' ', filename)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    
    words = jieba.cut(cleaned, HMM=False)
    filtered = []
    for w in words:
        if not w or w.strip() == '':
            continue
        if w not in STOPWORDS:
            if len(w) > 1 or w in MEANINGFUL_SINGLE_CHARS:
                filtered.append(w)
    
    if not filtered:
        # 如果全部被过滤，保留所有非空词
        words = jieba.cut(cleaned, HMM=False)
        filtered = [w for w in words if w.strip()]
    
    if not filtered:
        return cleaned
    
    result = " ".join(filtered)
    
    # 清空临时列表
    filtered.clear()
    del filtered
    
    return result


# ---------- 核心预测函数（内存优化）----------
def predict_pptx(filepath, verbose=True):
    """
    预测单个PPTX文件（优化内存，每次预测后释放中间变量）
    """
    global _model, _text_tokenizer, _filename_tokenizer, _categories, _config, _prediction_count
    
    if _model is None:
        load_models()
    
    filepath = Path(filepath)
    
    if not filepath.exists():
        raise FileNotFoundError(f"文件不存在: {filepath}")
    
    if not filepath.suffix.lower() in ['.pptx', '.ppt']:
        raise ValueError(f"不支持的文件格式: {filepath.suffix}")
    
    # 1. 提取特征
    filename_feature = extract_filename_features(str(filepath))
    raw_text = extract_text_from_pptx(str(filepath))
    
    if raw_text:
        cleaned = clean_text(raw_text)
        text_feature = cut_words(cleaned)
        text_available = True
    else:
        text_feature = "空文本"
        text_available = False
        if verbose:
            print(f"  ⚠ 警告: 未提取到文本内容，仅使用文件名分类")
    
    # 2. 转换为序列
    text_seq = _text_tokenizer.texts_to_sequences([text_feature])
    text_pad = pad_sequences(text_seq, maxlen=_config['max_sequence_length'], 
                            padding='post', truncating='post')
    
    filename_seq = _filename_tokenizer.texts_to_sequences([filename_feature])
    filename_pad = pad_sequences(filename_seq, maxlen=_config['max_filename_length'],
                                padding='post', truncating='post')
    
    # 3. 预测（使用静默模式）
    import io
    old_stdout = sys.stdout
    sys.stdout = io.StringIO()
    
    try:
        predictions = _model.predict([text_pad, filename_pad], verbose=0, batch_size=1)
    finally:
        sys.stdout = old_stdout
    
    pred_idx = predictions.argmax()
    confidence = float(predictions[0][pred_idx])
    predicted_class = _categories[pred_idx]
    
    # 获取所有类别的置信度（只取前5，节省内存）
    all_probs = {cat: float(predictions[0][i]) for i, cat in enumerate(_categories)}
    sorted_probs = sorted(all_probs.items(), key=lambda x: x[1], reverse=True)
    
    if verbose:
        print(f"\n{'='*50}")
        print(f"文件: {filepath.name}")
        print(f"{'='*50}")
        print(f"📄 文件名: {filename_feature[:80]}{'...' if len(filename_feature) > 80 else ''}")
        print(f"📝 文本长度: {len(text_feature.split())} 词")
        print(f"\n🎯 预测结果: {predicted_class}")
        print(f"📊 置信度: {confidence:.2%}")
        print(f"\n📈 详细分类概率:")
        for cat, prob in sorted_probs[:5]:
            bar_length = int(prob * 30)
            bar = "█" * bar_length if bar_length > 0 else "·"
            print(f"   {cat:4s}: {prob:6.2%} {bar}")
    
    # 主动释放中间变量（避免内存累积）
    del text_seq
    del text_pad
    del filename_seq
    del filename_pad
    del predictions
    del all_probs
    del sorted_probs
    
    # 每预测5次后强制垃圾回收
    _prediction_count += 1
    
    if _prediction_count % 5 == 0:
        gc.collect()
        try:
            tf.keras.backend.clear_session()
        except:
            pass
    
    return predicted_class, confidence, {
        'filename_feature': filename_feature,
        'text_length': len(text_feature.split()),
        'all_probabilities': None,
        'text_available': text_available
    }


# ---------- 交互式预测（内存优化版）----------
def interactive_mode():
    """交互式预测模式 - 内存优化版"""
    global _prediction_count
    
    print("\n" + "="*50)
    print("PPTX学科分类器 - 交互式预测模式（内存优化版）")
    print("="*50)
    print("输入PPTX文件路径进行预测，输入 'quit' 或 'exit' 退出")
    print("输入 'mem' 查看当前内存使用情况")
    print("提示：可以直接拖拽文件到命令行窗口")
    print("-"*50)
    
    # 重置预测计数器
    _prediction_count = 0
    
    # 尝试导入psutil（可选）
    psutil_available = False
    try:
        import psutil
        psutil_available = True
        process = psutil.Process()
        initial_memory = process.memory_info().rss / 1024 / 1024
        print(f"初始内存: {initial_memory:.1f} MB")
    except ImportError:
        print("（提示：安装psutil可查看内存使用情况: pip install psutil）")
    
    prediction_count = 0
    
    while True:
        try:
            filepath = input("\n请输入文件路径: ").strip().strip('"').strip("'")
        except (KeyboardInterrupt, EOFError):
            print("\n再见！")
            break
        
        if filepath.lower() in ['quit', 'exit', 'q', "bye"]:
            print("再见！")
            break
        
        if filepath.lower() == 'mem' and psutil_available:
            current_memory = process.memory_info().rss / 1024 / 1024
            print(f"当前内存: {current_memory:.1f} MB | 已预测次数: {prediction_count}")
            continue
        
        if not filepath:
            continue
        
        if not Path(filepath).exists():
            print(f"错误: 文件不存在 '{filepath}'")
            continue
        
        try:
            predict_pptx(filepath)
            prediction_count += 1
            
            # 每10次预测后显示内存状态
            if prediction_count % 10 == 0 and psutil_available:
                current_memory = process.memory_info().rss / 1024 / 1024
                print(f"\n📊 [内存状态] 当前: {current_memory:.1f} MB | 已预测: {prediction_count} 次")
                
                # 主动进行垃圾回收
                gc.collect()
                try:
                    tf.keras.backend.clear_session()
                except:
                    pass
                
        except Exception as e:
            print(f"预测失败: {e}")


# ---------- 批量预测 ----------
def predict_batch(directory_path, recursive=False):
    """批量预测（一次性处理，不累积内存）"""
    global _prediction_count
    _prediction_count = 0
    
    directory = Path(directory_path)
    
    if not directory.exists():
        raise FileNotFoundError(f"目录不存在: {directory}")
    
    if recursive:
        pptx_files = list(directory.rglob("*.pptx")) + list(directory.rglob("*.PPTX"))
    else:
        pptx_files = list(directory.glob("*.pptx")) + list(directory.glob("*.PPTX"))
    
    if not pptx_files:
        print(f"未找到PPTX文件: {directory}")
        return []
    
    print(f"\n找到 {len(pptx_files)} 个PPTX文件，开始批量预测...")
    print("="*60)
    
    results = []
    
    # 尝试导入psutil
    psutil_available = False
    try:
        import psutil
        psutil_available = True
        process = psutil.Process()
    except ImportError:
        pass
    
    for i, filepath in enumerate(pptx_files, 1):
        print(f"\n[{i}/{len(pptx_files)}] 处理中...", end=' ')
        try:
            pred_class, confidence, details = predict_pptx(str(filepath), verbose=False)
            results.append({
                'file': str(filepath),
                'predicted_class': pred_class,
                'confidence': confidence,
                'text_length': details['text_length'],
                'text_available': details['text_available']
            })
            print(f"✓ {filepath.name} → {pred_class} ({confidence:.1%})")
        except Exception as e:
            print(f"✗ {filepath.name} → 错误: {e}")
            results.append({
                'file': str(filepath),
                'predicted_class': 'ERROR',
                'confidence': 0,
                'error': str(e)
            })
        
        # 每处理10个文件清理一次内存
        if i % 10 == 0:
            gc.collect()
            if psutil_available:
                current_memory = process.memory_info().rss / 1024 / 1024
                print(f"\n  [内存: {current_memory:.1f} MB]")
    
    # 打印汇总
    print("\n" + "="*60)
    print("批量预测汇总")
    print("="*60)
    
    from collections import Counter
    class_counts = Counter([r['predicted_class'] for r in results if 'predicted_class' in r and r['predicted_class'] != 'ERROR'])
    
    if class_counts:
        print("\n类别分布:")
        for cat, count in sorted(class_counts.items()):
            print(f"  {cat}: {count} 个")
    
    # 导出结果
    output_csv = directory / f"prediction_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
    import csv
    with open(output_csv, 'w', newline='', encoding='utf-8-sig') as f:
        writer = csv.DictWriter(f, fieldnames=['file', 'predicted_class', 'confidence', 'text_length', 'text_available'])
        writer.writeheader()
        for r in results:
            if 'error' not in r:
                writer.writerow({k: r.get(k, '') for k in ['file', 'predicted_class', 'confidence', 'text_length', 'text_available']})
    
    print(f"\n结果已导出到: {output_csv}")
    
    return results


# ---------- 命令行入口 ----------
def main():
    # 加载模型
    print("正在初始化...", end=' ', flush=True)
    try:
        load_models()
    except Exception as e:
        print(f"\n错误: 模型加载失败 - {e}")
        print("\n请确保以下文件存在于当前目录:")
        for f in [MODEL_PATH, TEXT_TOKENIZER_PATH, FILENAME_TOKENIZER_PATH, CATEGORIES_PATH, CONFIG_PATH]:
            exists = "✓" if Path(f).exists() else "✗"
            print(f"  {exists} {f}")
        return
    
    parser = argparse.ArgumentParser(description='PPTX学科分类器预测工具（内存优化版）')
    parser.add_argument('input', nargs='?', help='PPTX文件路径或目录路径')
    parser.add_argument('--batch', '-b', action='store_true', help='批量预测模式')
    parser.add_argument('--recursive', '-r', action='store_true', help='递归搜索子目录')
    parser.add_argument('--interactive', '-i', action='store_true', help='交互式预测模式')
    
    args = parser.parse_args()
    
    if args.interactive:
        interactive_mode()
    elif args.batch or (args.input and Path(args.input).is_dir()):
        input_path = args.input or "."
        predict_batch(input_path, recursive=args.recursive)
    elif args.input:
        predict_pptx(args.input)
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
