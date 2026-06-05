#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX学科分类器 API服务 - 本地路径版本
提供RESTful API接口，接收文件路径进行预测
无文件上传功能，直接使用服务器本地路径
"""

import os
import sys
import gc
import warnings
import re
import pickle
import json
from pathlib import Path
from datetime import datetime
from collections import Counter
from flask import Flask, request, jsonify
from flask_cors import CORS
import numpy as np
import logging

# 屏蔽所有警告
warnings.filterwarnings('ignore')
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'

# ONNX Runtime
import onnxruntime as ort

import jieba
import jieba.analyse
from pptx import Presentation

# 创建Flask应用
app = Flask(__name__)
CORS(app)

# ---------- 配置 ----------
ONNX_MODEL_PATH = "textcnn_classifier.onnx"
TEXT_TOKENIZER_PATH = "text_tokenizer_none.pkl"
FILENAME_TOKENIZER_PATH = "filename_tokenizer_none.pkl"
CATEGORIES_PATH = "categories.pkl"
CONFIG_PATH = "config_optimized.pkl"

# 模型词汇表大小
MODEL_VOCAB_SIZE = 20000

# 关键词提取配置
KEYWORD_TOP_N = 12
KEYWORD_MIN_N = 5

# 停用词表
STOPWORDS = set([
    '的', '了', '是', '我', '你', '他', '她', '它', '我们', '你们', '他们',
    '这', '那', '有', '在', '不', '和', '与', '就', '都', '而', '及', '或',
    '一个', '这个', '那个', '那些', '这些', '这里', '那里', '然后', '因为',
    '所以', '但是', '如果', '虽然', '然而', '并且', '或者', '可以', '可能',
    '应该', '需要', '没有', '自己', '什么', '哪个', '如何', '为什么', '也',
    '还', '被', '把', '给', '让', '去', '年', '月', '日', '时', '分', '秒',
    '第', '一', '二', '三', '四', '五', '六', '七', '八', '九', '十',
    '上', '下', '中', '大', '小', '多', '少', '高', '低', '长', '短'
])

# 有意义的单字词
MEANINGFUL_SINGLE_CHARS = {
    '圆', '力', '氧', '氢', '碳', '钠', '酸', '碱', '盐', 
    '电', '光', '声', '热', '诗', '词', '歌', '曲', '数',
    '方', '程', '函', '数', '角', '形', '体', '积'
}

# 学科相关的高权重词
SUBJECT_KEYWORDS = {
    '语文': ['古诗', '文言文', '散文', '小说', '诗歌', '作者', '阅读', '写作', '作文', '修辞', '语法', '汉字', '成语', '唐诗', '宋词'],
    '数学': ['函数', '方程', '几何', '代数', '三角', '数列', '概率', '统计', '向量', '导数', '积分', '公式', '定理', '证明', '计算'],
    '英语': ['单词', '语法', '句型', '阅读', '听力', '写作', '翻译', '时态', '从句', '词汇', '发音', '对话', '短文', '字母'],
    '物理': ['力', '运动', '能量', '电场', '磁场', '电路', '光学', '热学', '量子', '相对论', '速度', '加速度', '质量', '密度', '压强'],
    '化学': ['反应', '元素', '化合物', '方程式', '原子', '分子', '离子', '酸碱', '氧化', '还原', '实验', '试剂', '催化剂', '溶液'],
    '生物': ['细胞', '基因', '生物', '植物', '动物', '生态', '进化', '遗传', '细菌', '病毒', '蛋白质', '酶', '光合作用', '呼吸作用'],
    '班会': ['主题', '活动', '班级', '同学', '老师', '纪律', '安全', '卫生', '德育', '心理健康', '团队', '合作', '文明', '礼仪']
}

# 抑制jieba输出
jieba.setLogLevel(logging.ERROR)
logging.getLogger().setLevel(logging.ERROR)

# ---------- 全局变量 ----------
_session = None
_text_tokenizer = None
_filename_tokenizer = None
_categories = None
_config = None
_prediction_count = 0


# ---------- 轻量级Tokenizer ----------
class SimpleTokenizer:
    def __init__(self, word_index, vocab_size, max_vocab_size=MODEL_VOCAB_SIZE):
        self.original_word_index = word_index
        self.max_vocab_size = max_vocab_size
        self.vocab_size = vocab_size
        
        self.word_index = {}
        self.oov_token = 1
        
        valid_words = []
        for word, idx in word_index.items():
            if 1 <= idx < max_vocab_size:
                valid_words.append((word, idx))
        
        valid_words.sort(key=lambda x: x[1])
        
        for new_idx, (word, old_idx) in enumerate(valid_words, 1):
            self.word_index[word] = new_idx
        
        self.compressed_size = len(self.word_index)
        self.index_word = {v: k for k, v in self.word_index.items()}
        
        print(f"  Tokenizer压缩: 原始={len(word_index)} -> 压缩后={self.compressed_size}")
    
    def texts_to_sequences(self, texts):
        sequences = []
        for text in texts:
            if isinstance(text, str):
                words = text.split()
                seq = []
                for word in words:
                    idx = self.word_index.get(word, self.oov_token)
                    if idx >= self.max_vocab_size:
                        idx = self.oov_token
                    seq.append(idx)
                sequences.append(seq)
            else:
                sequences.append([])
        return sequences
    
    def get_vocab_info(self):
        return {
            'original_size': len(self.original_word_index),
            'compressed_size': self.compressed_size,
            'max_supported': self.max_vocab_size - 1,
            'oov_token': self.oov_token
        }


def load_tokenizer_compressed(filepath, max_vocab_size=MODEL_VOCAB_SIZE):
    with open(filepath, 'rb') as f:
        data = pickle.load(f)
    
    if isinstance(data, dict) and 'word_index' in data:
        word_index = data['word_index']
    elif hasattr(data, 'word_index'):
        word_index = data.word_index
    else:
        raise ValueError(f"不支持的tokenizer格式: {type(data)}")
    
    print(f"  📊 加载原始词表: {len(word_index)} 个词")
    tokenizer = SimpleTokenizer(word_index, len(word_index), max_vocab_size)
    return tokenizer


# ---------- ONNX Runtime配置 ----------
def get_onnx_session():
    global _session
    
    if _session is not None:
        return _session
    
    if not Path(ONNX_MODEL_PATH).exists():
        raise FileNotFoundError(f"ONNX模型文件不存在: {ONNX_MODEL_PATH}")
    
    sess_options = ort.SessionOptions()
    sess_options.enable_cpu_mem_arena = True
    sess_options.execution_mode = ort.ExecutionMode.ORT_SEQUENTIAL
    sess_options.graph_optimization_level = ort.GraphOptimizationLevel.ORT_ENABLE_ALL
    sess_options.intra_op_num_threads = 1
    sess_options.inter_op_num_threads = 1
    
    ort.set_default_logger_severity(3)
    
    available_providers = ort.get_available_providers()
    print(f"可用的执行提供者: {available_providers}")
    
    providers = ['CPUExecutionProvider']
    if 'CUDAExecutionProvider' in available_providers:
        providers.append('CUDAExecutionProvider')
    
    try:
        _session = ort.InferenceSession(
            ONNX_MODEL_PATH, 
            sess_options=sess_options,
            providers=providers
        )
        print("✅ ONNX模型加载成功！")
    except Exception as e:
        print(f"❌ 模型加载失败: {e}")
        raise
    
    return _session


# ---------- 关键词提取函数 ----------
def extract_keywords(text, top_n=KEYWORD_TOP_N, min_n=KEYWORD_MIN_N):
    if not text or len(text.strip()) < 10:
        return ["无足够文本内容"]
    
    keywords_set = set()
    
    try:
        tfidf_keywords = jieba.analyse.extract_tags(
            text, topK=top_n, withWeight=False, allowPOS=()
        )
        keywords_set.update(tfidf_keywords[:top_n//2])
    except:
        pass
    
    try:
        textrank_keywords = jieba.analyse.textrank(
            text, topK=top_n, withWeight=False, allowPOS=('ns', 'n', 'vn', 'v', 'an')
        )
        keywords_set.update(textrank_keywords[:top_n//2])
    except:
        pass
    
    words = jieba.cut(text, HMM=False)
    word_freq = Counter()
    for w in words:
        w = w.strip()
        if len(w) >= 2 and w not in STOPWORDS and not w.isdigit():
            word_freq[w] += 1
    
    freq_keywords = [w for w, _ in word_freq.most_common(top_n)]
    keywords_set.update(freq_keywords[:top_n//2])
    
    for category, subject_words in SUBJECT_KEYWORDS.items():
        for sw in subject_words:
            if sw in text and len(sw) >= 2:
                keywords_set.add(sw)
    
    keywords = list(keywords_set)
    
    def keyword_score(kw):
        score = 0
        if 2 <= len(kw) <= 4:
            score += 10
        score += text.count(kw) * 5
        for subject_words in SUBJECT_KEYWORDS.values():
            if kw in subject_words:
                score += 20
                break
        return score
    
    keywords.sort(key=keyword_score, reverse=True)
    
    if len(keywords) < min_n:
        common_words = [w for w, _ in word_freq.most_common(top_n * 2) 
                       if w not in keywords and len(w) >= 2]
        keywords.extend(common_words[:min_n - len(keywords)])
    
    return keywords[:top_n]


def extract_filename_keywords(filename_text, top_n=5):
    if not filename_text:
        return []
    
    words = jieba.cut(filename_text, HMM=False)
    keywords = []
    for w in words:
        w = w.strip()
        if len(w) >= 2 and w not in STOPWORDS and not w.isdigit():
            keywords.append(w)
    
    seen = set()
    unique_keywords = []
    for kw in keywords:
        if kw not in seen:
            seen.add(kw)
            unique_keywords.append(kw)
    
    return unique_keywords[:top_n]


# ---------- 模型加载 ----------
def load_models():
    global _text_tokenizer, _filename_tokenizer, _categories, _config, _session
    
    required_files = [
        (ONNX_MODEL_PATH, "ONNX模型"),
        (TEXT_TOKENIZER_PATH, "文本分词器"),
        (FILENAME_TOKENIZER_PATH, "文件名词典"),
        (CATEGORIES_PATH, "类别映射"),
        (CONFIG_PATH, "配置")
    ]
    
    for filepath, desc in required_files:
        if not Path(filepath).exists():
            raise FileNotFoundError(f"{desc}文件不存在: {filepath}")
    
    print("\n📚 加载模型组件...")
    print("-" * 40)
    
    print("加载文本分词器...", flush=True)
    _text_tokenizer = load_tokenizer_compressed(TEXT_TOKENIZER_PATH, MODEL_VOCAB_SIZE)
    
    print("加载文件名词典...", flush=True)
    _filename_tokenizer = load_tokenizer_compressed(FILENAME_TOKENIZER_PATH, MODEL_VOCAB_SIZE)
    
    print("加载类别映射...", end=' ', flush=True)
    with open(CATEGORIES_PATH, 'rb') as f:
        _categories = pickle.load(f)
    print("✓")
    
    print("加载配置...", end=' ', flush=True)
    with open(CONFIG_PATH, 'rb') as f:
        _config = pickle.load(f)
    print("✓")
    
    print("-" * 40)
    print(f"支持 {len(_categories)} 个类别: {', '.join(_categories)}")
    
    _session = get_onnx_session()
    
    text_info = _text_tokenizer.get_vocab_info()
    filename_info = _filename_tokenizer.get_vocab_info()
    
    print(f"\n词表信息汇总:")
    print(f"  文本词表: {text_info['compressed_size']}/{text_info['original_size']} 个词")
    print(f"  文件名词表: {filename_info['compressed_size']}/{filename_info['original_size']} 个词")
    print()


def pad_sequences(sequences, maxlen, padding='post', truncating='post'):
    if not sequences:
        return np.zeros((0, maxlen), dtype=np.int32)
    
    result = np.zeros((len(sequences), maxlen), dtype=np.int32)
    
    for i, seq in enumerate(sequences):
        if truncating == 'post':
            seq = seq[:maxlen]
        else:
            seq = seq[-maxlen:] if len(seq) > maxlen else seq
        
        if padding == 'post':
            result[i, :len(seq)] = seq
        else:
            if len(seq) > 0:
                result[i, -len(seq):] = seq
    
    return result


# ---------- PPTX 文本提取 ----------
def extract_text_from_pptx(pptx_path):
    text_parts = []
    prs = None
    
    try:
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text.strip())
                    
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    text_parts.append(cell.text.strip())
                    
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text:
                                text_parts.append(paragraph.text.strip())
                except:
                    continue
            
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        text_parts.append(notes_slide.notes_text_frame.text.strip())
            except:
                pass
            
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    text_parts.append(slide.shapes.title.text.strip())
            except:
                pass
        
        result = " ".join(text_parts)
        
    except Exception as e:
        result = ""
    finally:
        if prs is not None:
            del prs
        text_parts.clear()
    
    return result


# ---------- 文本预处理 ----------
def clean_text(text):
    if not text:
        return ""
    
    text = re.sub(r'http\S+|www\S+|https\S+', '', text)
    text = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9，。！？；：""''（）【】《》、 ]', '', text)
    text = re.sub(r'\s+', ' ', text)
    
    return text.strip()


def cut_words(text):
    if not text:
        return ""
    
    words = jieba.cut(text, HMM=False)
    filtered = []
    for w in words:
        if not w or w.strip() == '':
            continue
        if w not in STOPWORDS:
            filtered.append(w)
    
    if not filtered:
        words = jieba.cut(text, HMM=False)
        filtered = [w for w in words if w.strip()]
    
    if not filtered:
        return "空文本"
    
    return " ".join(filtered)


def extract_filename_features(filepath):
    filepath_obj = Path(filepath) if isinstance(filepath, str) else filepath
    filename = filepath_obj.stem
    
    cleaned = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', ' ', filename)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    
    words = jieba.cut(cleaned, HMM=False)
    filtered = []
    for w in words:
        if not w or w.strip() == '':
            continue
        if w not in STOPWORDS:
            if len(w) > 1 or w in MEANINGFUL_SINGLE_CHARS:
                filtered.append(w)
    
    if not filtered:
        words = jieba.cut(cleaned, HMM=False)
        filtered = [w for w in words if w.strip()]
    
    if not filtered:
        return cleaned
    
    return " ".join(filtered)


# ---------- 核心预测函数（接收路径）----------
def predict_pptx_by_path(filepath):
    """
    预测函数，直接接收文件路径
    """
    global _session, _text_tokenizer, _filename_tokenizer, _categories, _config, _prediction_count
    
    filepath = Path(filepath)
    
    if not filepath.exists():
        raise FileNotFoundError(f"文件不存在: {filepath}")
    
    # 提取特征
    filename_feature = extract_filename_features(str(filepath))
    raw_text = extract_text_from_pptx(str(filepath))
    
    if raw_text:
        cleaned = clean_text(raw_text)
        text_feature = cut_words(cleaned)
        text_available = True
        keywords = extract_keywords(raw_text)
        filename_keywords = extract_filename_keywords(filename_feature)
    else:
        text_feature = "空文本"
        text_available = False
        keywords = ["无文本内容"]
        filename_keywords = extract_filename_keywords(filename_feature)
    
    # 转换为序列
    text_seq = _text_tokenizer.texts_to_sequences([text_feature])
    text_pad = pad_sequences(
        text_seq, 
        maxlen=_config.get('max_sequence_length', 750),
        padding='post', 
        truncating='post'
    )
    
    filename_seq = _filename_tokenizer.texts_to_sequences([filename_feature])
    filename_pad = pad_sequences(
        filename_seq, 
        maxlen=_config.get('max_filename_length', 30),
        padding='post', 
        truncating='post'
    )
    
    text_pad = text_pad.astype(np.int32)
    filename_pad = filename_pad.astype(np.int32)
    
    # ONNX推理
    input_names = [inp.name for inp in _session.get_inputs()]
    output_names = [out.name for out in _session.get_outputs()]
    
    inputs = {
        input_names[0]: text_pad,
        input_names[1]: filename_pad
    }
    
    predictions = _session.run(output_names, inputs)[0]
    
    pred_idx = int(predictions[0].argmax())
    confidence = float(predictions[0][pred_idx])
    predicted_class = _categories[pred_idx]
    
    all_probabilities = {
        cat: float(predictions[0][i]) 
        for i, cat in enumerate(_categories)
    }
    
    _prediction_count += 1
    if _prediction_count % 5 == 0:
        gc.collect()
    
    return {
        'filename': filepath.name,
        'filepath': str(filepath),
        'predicted_class': predicted_class,
        'confidence': confidence,
        'all_probabilities': all_probabilities,
        'text_available': text_available,
        'text_length': len(text_feature.split()),
        'filename_features': filename_feature,
        'keywords': keywords,
        'filename_keywords': filename_keywords
    }


# ---------- API 路由 ----------
@app.route('/health', methods=['GET'])
def health_check():
    """健康检查接口"""
    return jsonify({
        'status': 'healthy',
        'model_loaded': _session is not None,
        'categories': _categories,
        'backend': 'onnxruntime',
        'timestamp': datetime.now().isoformat()
    })


@app.route('/predict', methods=['POST'])
def predict_single():
    """
    单文件预测接口（接收文件路径）
    请求：POST JSON格式
    {
        "filepath": "/path/to/your/file.pptx"
    }
    """
    if _session is None:
        return jsonify({'error': '模型未加载'}), 503
    
    data = request.get_json()
    if not data:
        return jsonify({'error': '请求体必须是JSON格式'}), 400
    
    filepath = data.get('filepath')
    if not filepath:
        return jsonify({'error': '请提供文件路径', 'required_field': 'filepath'}), 400
    
    # 支持相对路径转换为绝对路径
    filepath = os.path.abspath(filepath)
    
    # 检查文件扩展名
    if not filepath.lower().endswith(('.pptx', '.ppt')):
        return jsonify({'error': '不支持的文件格式，请提供PPTX或PPT文件'}), 400
    
    try:
        result = predict_pptx_by_path(filepath)
        
        result['success'] = True
        result['top_3'] = sorted(
            result['all_probabilities'].items(), 
            key=lambda x: x[1], reverse=True
        )[:3]
        result['keyword_count'] = len(result['keywords'])
        
        return jsonify(result)
        
    except FileNotFoundError as e:
        return jsonify({'error': str(e)}), 404
    except Exception as e:
        return jsonify({'error': f'预测失败: {str(e)}'}), 500


@app.route('/predict/batch', methods=['POST'])
def predict_batch():
    """
    批量预测接口
    请求：POST JSON格式
    {
        "filepaths": ["/path/to/file1.pptx", "/path/to/file2.pptx"]
    }
    """
    if _session is None:
        return jsonify({'error': '模型未加载'}), 503
    
    data = request.get_json()
    if not data:
        return jsonify({'error': '请求体必须是JSON格式'}), 400
    
    filepaths = data.get('filepaths', [])
    if not filepaths or len(filepaths) == 0:
        return jsonify({'error': '请提供文件路径列表', 'required_field': 'filepaths'}), 400
    
    results = []
    errors = []
    
    for filepath in filepaths:
        filepath = os.path.abspath(filepath)
        
        if not filepath.lower().endswith(('.pptx', '.ppt')):
            errors.append({
                'success': False,
                'filepath': filepath,
                'error': '不支持的文件格式'
            })
            continue
        
        try:
            result = predict_pptx_by_path(filepath)
            results.append({
                'success': True,
                'filepath': result['filepath'],
                'filename': result['filename'],
                'predicted_class': result['predicted_class'],
                'confidence': result['confidence'],
                'text_available': result['text_available'],
                'text_length': result['text_length'],
                'keywords': result['keywords'][:5],
                'filename_keywords': result['filename_keywords']
            })
        except FileNotFoundError as e:
            errors.append({
                'success': False,
                'filepath': filepath,
                'error': str(e)
            })
        except Exception as e:
            errors.append({
                'success': False,
                'filepath': filepath,
                'error': f'预测失败: {str(e)}'
            })
    
    # 统计信息
    summary = {}
    for r in results:
        if r['success']:
            cat = r['predicted_class']
            summary[cat] = summary.get(cat, 0) + 1
    
    return jsonify({
        'success': True,
        'total': len(filepaths),
        'success_count': len(results),
        'error_count': len(errors),
        'summary': summary,
        'results': results,
        'errors': errors,
        'timestamp': datetime.now().isoformat(),
        'backend': 'onnxruntime'
    })


@app.route('/categories', methods=['GET'])
def get_categories():
    """获取所有支持的类别"""
    return jsonify({
        'categories': _categories,
        'count': len(_categories),
        'timestamp': datetime.now().isoformat()
    })


@app.route('/stats', methods=['GET'])
def get_stats():
    """获取服务统计信息"""
    try:
        import psutil
        process = psutil.Process()
        memory_info = process.memory_info()
        memory_mb = memory_info.rss / 1024 / 1024
    except:
        memory_mb = -1
    
    return jsonify({
        'predictions_count': _prediction_count,
        'memory_usage_mb': memory_mb,
        'model_loaded': _session is not None,
        'categories': _categories,
        'backend': 'onnxruntime',
        'keyword_config': {
            'top_n': KEYWORD_TOP_N,
            'min_n': KEYWORD_MIN_N
        },
        'timestamp': datetime.now().isoformat()
    })


@app.errorhandler(404)
def not_found(error):
    return jsonify({'error': '接口不存在'}), 404


@app.errorhandler(500)
def internal_error(error):
    return jsonify({'error': '服务器内部错误'}), 500


# ---------- 启动服务 ----------
def main():
    global _prediction_count
    _prediction_count = 0
    
    print("="*50)
    print("PPTX学科分类器 API服务启动 - 本地路径版本")
    print("="*50)
    
    try:
        load_models()
    except Exception as e:
        print(f"\n❌ 错误: 模型加载失败 - {e}")
        print("\n请确保以下文件存在于当前目录:")
        for f in [ONNX_MODEL_PATH, TEXT_TOKENIZER_PATH, FILENAME_TOKENIZER_PATH, CATEGORIES_PATH, CONFIG_PATH]:
            exists = "✓" if Path(f).exists() else "✗"
            print(f"  {exists} {f}")
        sys.exit(1)
    
    port = int(os.environ.get('PORT', 5000))
    host = os.environ.get('HOST', '0.0.0.0')
    debug = os.environ.get('DEBUG', 'False').lower() == 'true'
    
    print(f"\n服务启动信息:")
    print(f"  地址: http://{host}:{port}")
    print(f"  调试模式: {debug}")
    print(f"  推理后端: ONNX Runtime")
    print(f"\n可用接口:")
    print("  POST /predict          - 单文件预测（JSON: {\"filepath\": \"路径\"}）")
    print(f"  POST /predict/batch    - 批量预测（JSON: {{\"filepaths\": [\"路径1\", \"路径2\"]}}）")
    print(f"  GET  /health           - 健康检查")
    print(f"  GET  /categories       - 获取类别列表")
    print(f"  GET  /stats            - 获取统计信息")
    print("="*50)
    
    app.run(host=host, port=port, debug=debug, threaded=True)


if __name__ == "__main__":
    main()