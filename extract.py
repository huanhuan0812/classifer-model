#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件文字提取器 - 纯提取版本（不进行任何清洗）
保留原始文本中的所有内容
"""

import os
import gc
import zipfile
from pathlib import Path
from io import BytesIO

try:
    from pptx import Presentation
except ImportError:
    raise ImportError("请安装 python-pptx: pip install python-pptx")

try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("警告: python-docx未安装，DOCX支持将受限")


def extract_text_from_docx(docx_path):
    """
    从DOCX文件中提取所有文本内容（不进行清洗）
    
    Args:
        docx_path: DOCX文件路径
    
    Returns:
        str: 原始文本内容
    """
    all_text = []
    
    if not DOCX_SUPPORT:
        return ""
    
    try:
        doc = Document(docx_path)
        
        # 提取段落文本
        for paragraph in doc.paragraphs:
            if paragraph.text:
                all_text.append(paragraph.text)
        
        # 提取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text:
                        row_text.append(cell.text.strip())
                if row_text:
                    all_text.append(" | ".join(row_text))
        
        # 提取页眉
        if hasattr(doc, 'sections'):
            for section in doc.sections:
                if section.header and section.header.paragraphs:
                    header_text = []
                    for para in section.header.paragraphs:
                        if para.text:
                            header_text.append(para.text)
                    if header_text:
                        all_text.append("[页眉] " + " ".join(header_text))
                
                if section.footer and section.footer.paragraphs:
                    footer_text = []
                    for para in section.footer.paragraphs:
                        if para.text:
                            footer_text.append(para.text)
                    if footer_text:
                        all_text.append("[页脚] " + " ".join(footer_text))
        
        return "\n".join(all_text) if all_text else ""
        
    except Exception as e:
        return ""
    finally:
        gc.collect()


def extract_embedded_docx_from_pptx(pptx_path):
    """
    从PPTX文件中提取嵌入的DOCX文件内容
    
    Args:
        pptx_path: PPTX文件路径
    
    Returns:
        str: 嵌入DOCX的文本内容
    """
    embedded_texts = []
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for file_info in zf.filelist:
                if file_info.filename.startswith('ppt/embeddings/') and file_info.filename.endswith(('.docx', '.doc')):
                    try:
                        docx_data = zf.read(file_info.filename)
                        
                        if DOCX_SUPPORT:
                            docx_stream = BytesIO(docx_data)
                            doc = Document(docx_stream)
                            
                            for paragraph in doc.paragraphs:
                                if paragraph.text:
                                    embedded_texts.append(paragraph.text)
                            
                            for table in doc.tables:
                                for row in table.rows:
                                    row_text = []
                                    for cell in row.cells:
                                        if cell.text:
                                            row_text.append(cell.text.strip())
                                    if row_text:
                                        embedded_texts.append("[嵌入表格] " + " | ".join(row_text))
                            
                            del doc
                        else:
                            # 备用方法：直接解析XML
                            try:
                                docx_zip = zipfile.ZipFile(BytesIO(docx_data))
                                if 'word/document.xml' in docx_zip.namelist():
                                    import xml.etree.ElementTree as ET
                                    xml_content = docx_zip.read('word/document.xml')
                                    root = ET.fromstring(xml_content)
                                    namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                                    
                                    for t in root.findall('.//w:t', namespaces):
                                        if t.text:
                                            embedded_texts.append(t.text)
                                    
                                docx_zip.close()
                            except:
                                pass
                    except:
                        pass
    except:
        pass
    
    return "\n".join(embedded_texts) if embedded_texts else ""


def extract_text_from_pptx(pptx_path, extract_embedded=True):
    """
    从PPTX文件中提取所有文本内容（不进行清洗）
    
    Args:
        pptx_path: PPTX文件路径
        extract_embedded: 是否提取嵌入的DOCX文件
    
    Returns:
        str: 原始文本内容
    """
    all_text = []
    prs = None
    
    try:
        # 提取嵌入的docx文件内容
        if extract_embedded:
            embedded_text = extract_embedded_docx_from_pptx(pptx_path)
            if embedded_text:
                all_text.append("[嵌入DOCX内容]")
                all_text.append(embedded_text)
                all_text.append("")
        
        prs = Presentation(pptx_path)
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"[幻灯片 {slide_idx}]")
            
            try:
                # 提取标题
                if slide.shapes.title and slide.shapes.title.text:
                    slide_text.append(f"标题: {slide.shapes.title.text}")
            except:
                pass
            
            # 提取所有形状中的文本
            for shape in slide.shapes:
                try:
                    # 提取形状文本
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # 提取表格中的文本
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text)
                            if row_text:
                                slide_text.append("[表格] " + " | ".join(row_text))
                    
                    # 检查文本框
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text_frame.text:
                            slide_text.append(shape.text_frame.text)
                            
                except:
                    continue
            
            # 提取备注
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        slide_text.append("[备注] " + notes_slide.notes_text_frame.text)
            except:
                pass
            
            if slide_text:
                all_text.append("\n".join(slide_text))
                all_text.append("")
            
            del slide
        
    except Exception as e:
        return ""
    finally:
        if prs is not None:
            del prs
        gc.collect()
    
    return "\n".join(all_text) if all_text else ""


def extract_text_from_file(filepath, extract_embedded=True):
    """
    根据文件类型提取文本内容（主入口函数，不进行清洗）
    
    Args:
        filepath: 文件路径 (str 或 Path)
        extract_embedded: 是否从PPTX中提取嵌入的DOCX内容
    
    Returns:
        str: 原始文本内容
    
    Raises:
        ValueError: 不支持的文件格式
        FileNotFoundError: 文件不存在
    """
    filepath = Path(filepath)
    
    if not filepath.exists():
        raise FileNotFoundError(f"文件不存在: {filepath}")
    
    suffix = filepath.suffix.lower()
    
    if suffix in ['.pptx', '.ppt']:
        return extract_text_from_pptx(str(filepath), extract_embedded)
    elif suffix == '.docx':
        return extract_text_from_docx(str(filepath))
    else:
        raise ValueError(f"不支持的文件格式: {suffix}，仅支持 .pptx, .ppt, .docx")


# ============ 命令行使用 ============
if __name__ == "__main__":
    import sys
    
    print("=" * 60)
    print("文件文字提取器（纯提取版本 - 不清洗）")
    print("支持格式: PPTX, PPT, DOCX")
    print(f"DOCX支持: {'✅ 已启用' if DOCX_SUPPORT else '❌ 未安装'}")
    print("=" * 60)
    
    if len(sys.argv) < 2:
        print("\n用法: python extract_raw.py <文件路径>")
        print("示例: python extract_raw.py example.pptx")
        sys.exit(0)
    
    filepath = sys.argv[1]
    
    try:
        print(f"\n📄 正在提取: {filepath}")
        print("-" * 60)
        
        # 提取文本（不清洗）
        text = extract_text_from_file(filepath)
        
        # 输出统计信息
        print(f"📝 提取字符数: {len(text)}")
        print(f"📝 行数: {len(text.splitlines()) if text else 0}")
        
        print("\n" + "=" * 60)
        print("提取的原始文本内容:")
        print("=" * 60)
        print(text if text else "[未提取到任何文本]")
        
        # 保存到文件
        output_path = Path(filepath).with_suffix('.txt')
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"\n✅ 文本已保存到: {output_path}")
        
    except Exception as e:
        print(f"❌ 提取失败: {e}")
        import traceback
        traceback.print_exc()