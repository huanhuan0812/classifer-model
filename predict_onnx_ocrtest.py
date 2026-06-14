#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX学科分类器预测脚本 - ONNX版本（零缓存，无Keras依赖）
支持7类别：语文、数学、英语、物理、化学、生物、班会
支持文件格式：.pptx, .ppt, .docx
使用ONNX Runtime进行推理，内存占用更小，速度更快
支持Tesseract OCR：当提取文本内容较少时自动识别PPT中的图片文字
"""

import os
import sys
import gc
import warnings
import re
import pickle
import argparse
import logging
import zipfile
import tempfile
from pathlib import Path
from datetime import datetime
from io import BytesIO

# 屏蔽TensorFlow和Keras相关警告（如果有）
os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3'
os.environ['TF_ENABLE_ONEDNN_OPTS'] = '0'

# 屏蔽所有警告
warnings.filterwarnings('ignore')

# 设置logging级别
logging.getLogger().setLevel(logging.ERROR)

import numpy as np
import onnxruntime as ort
import jieba
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 尝试导入docx处理库
try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    print("警告: python-docx未安装，docx文件支持将受限。请运行: pip install python-docx")

# 尝试导入OCR相关库
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
    # 尝试设置tesseract路径（可选）
    if sys.platform == 'win32':
        # Windows常见安装路径
        possible_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
        ]
        for path in possible_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                break
    print("✅ Tesseract OCR已启用（中英文识别）")
except ImportError:
    OCR_AVAILABLE = False
    print("提示: Tesseract OCR未安装，图片文字识别功能不可用。如需使用请运行:")
    print("  pip install pytesseract pillow")
    print("  并安装Tesseract-OCR引擎: https://github.com/tesseract-ocr/tesseract")

# ---------- 配置 ----------
ONNX_MODEL_PATH = "textcnn_classifier.onnx"
TEXT_TOKENIZER_PATH = "text_tokenizer_none.pkl"
FILENAME_TOKENIZER_PATH = "filename_tokenizer_none.pkl"
CATEGORIES_PATH = "categories.pkl"
CONFIG_PATH = "config_optimized.pkl"

# 模型词汇表大小（从错误信息确定：范围[-20000,19999]）
# 有效索引范围: 0(Padding), 1-19999(词汇)
MODEL_VOCAB_SIZE = 20000

# OCR启用阈值：当提取的文本字符数少于该值时启用OCR
OCR_MIN_TEXT_THRESHOLD = 100

# 停用词表
STOPWORDS = set([
    '的', '了', '是', '我', '你', '他', '她', '它', '我们', '你们', '他们',
    '这', '那', '有', '在', '不', '和', '与', '就', '都', '而', '及', '或',
    '一个', '这个', '那个', '那些', '这些', '这里', '那里', '然后', '因为',
    '所以', '但是', '如果', '虽然', '然而', '并且', '或者'
])

# 有意义的单字词
MEANINGFUL_SINGLE_CHARS = {'圆', '力', '氧', '氢', '碳', '钠', '酸', '碱', '盐', 
                           '电', '光', '声', '热', '诗', '词', '歌', '曲', '数',
                           '方', '程', '函', '数', '角', '形', '体', '积'}

# 抑制jieba输出
jieba.setLogLevel(logging.ERROR)

# ---------- 全局变量 ----------
_session = None
_text_tokenizer = None
_filename_tokenizer = None
_categories = None
_config = None


# ---------- 轻量级Tokenizer ----------
class SimpleTokenizer:
    """轻量级Tokenizer，强制限制索引在模型词汇表范围内"""
    def __init__(self, word_index, vocab_size, max_vocab_size=MODEL_VOCAB_SIZE):
        self.original_word_index = word_index
        self.max_vocab_size = max_vocab_size
        self.vocab_size = vocab_size
        
        # 构建压缩后的word_index（只保留模型支持的索引）
        self.word_index = {}
        self.oov_token = 1  # 使用索引1作为OOV
        
        # 过滤并重新映射
        valid_words = []
        for word, idx in word_index.items():
            if 1 <= idx < max_vocab_size:
                valid_words.append((word, idx))
        
        # 按原始索引排序
        valid_words.sort(key=lambda x: x[1])
        
        # 重新映射到连续索引（从1开始）
        for new_idx, (word, old_idx) in enumerate(valid_words, 1):
            self.word_index[word] = new_idx
        
        self.compressed_size = len(self.word_index)
        
        # 确保OOV token存在
        if self.oov_token not in self.word_index.values():
            # 如果索引1未被使用，保留它作为OOV
            pass
        
        self.index_word = {v: k for k, v in self.word_index.items()}
        
        print(f"  Tokenizer压缩: 原始={len(word_index)} -> 压缩后={self.compressed_size} (限制:{max_vocab_size-1})")
        print(f"     保留率: {self.compressed_size/len(word_index)*100:.1f}%")
    
    def texts_to_sequences(self, texts):
        """将文本转换为序列，所有索引都在模型有效范围内"""
        sequences = []
        for text in texts:
            if isinstance(text, str):
                words = text.split()
                seq = []
                for word in words:
                    # 获取索引，如果不在压缩词表中则使用OOV token
                    idx = self.word_index.get(word, self.oov_token)
                    # 双重保险：确保索引在有效范围内
                    if idx >= self.max_vocab_size:
                        idx = self.oov_token
                    seq.append(idx)
                sequences.append(seq)
            else:
                sequences.append([])
        return sequences
    
    def get_vocab_info(self):
        """获取词汇表信息"""
        return {
            'original_size': len(self.original_word_index),
            'compressed_size': self.compressed_size,
            'max_supported': self.max_vocab_size - 1,
            'oov_token': self.oov_token
        }


# ---------- 加载tokenizer（强制压缩到模型词汇表大小）----------
def load_tokenizer_compressed(filepath, max_vocab_size=MODEL_VOCAB_SIZE):
    """加载tokenizer并强制压缩词汇表到模型支持的大小"""
    with open(filepath, 'rb') as f:
        data = pickle.load(f)
    
    # 获取原始word_index
    if isinstance(data, dict) and 'word_index' in data:
        word_index = data['word_index']
    elif hasattr(data, 'word_index'):
        word_index = data.word_index
    else:
        raise ValueError(f"不支持的tokenizer格式: {type(data)}")
    
    print(f"  📊 加载原始词表: {len(word_index)} 个词")
    
    # 显示原始索引范围
    if word_index:
        min_idx = min(word_index.values())
        max_idx = max(word_index.values())
        print(f"     原始索引范围: {min_idx} - {max_idx}")
    
    # 创建压缩后的tokenizer
    tokenizer = SimpleTokenizer(word_index, len(word_index), max_vocab_size)
    
    return tokenizer


# ---------- ONNX Runtime配置 ----------
def get_onnx_session():
    """获取ONNX Runtime推理会话"""
    global _session
    
    if _session is not None:
        return _session
    
    if not Path(ONNX_MODEL_PATH).exists():
        raise FileNotFoundError(f"ONNX模型文件不存在: {ONNX_MODEL_PATH}")
    
    sess_options = ort.SessionOptions()
    sess_options.enable_cpu_mem_arena = True
    sess_options.execution_mode = ort.ExecutionMode.ORT_SEQUENTIAL
    sess_options.graph_optimization_level = ort.GraphOptimizationLevel.ORT_ENABLE_ALL
    sess_options.intra_op_num_threads = 1
    sess_options.inter_op_num_threads = 1
    sess_options.enable_profiling = False
    
    ort.set_default_logger_severity(3)
    
    available_providers = ort.get_available_providers()
    print(f"可用的执行提供者: {available_providers}")
    
    # 优先使用CPU，避免GPU内存问题
    providers = ['CPUExecutionProvider']
    if 'CUDAExecutionProvider' in available_providers:
        providers.append('CUDAExecutionProvider')
    
    try:
        _session = ort.InferenceSession(
            ONNX_MODEL_PATH, 
            sess_options=sess_options,
            providers=providers
        )
        print("✅ ONNX模型加载成功！")
        
        print(f"\n模型输入:")
        for inp in _session.get_inputs():
            shape_str = str(inp.shape).replace('None', 'batch')
            print(f"  - {inp.name}: {shape_str} ({inp.type})")
        
        print(f"\n模型输出:")
        for out in _session.get_outputs():
            shape_str = str(out.shape).replace('None', 'batch')
            print(f"  - {out.name}: {shape_str} ({out.type})")
        
        print(f"\n模型词汇表限制: {MODEL_VOCAB_SIZE} (有效索引: 1-{MODEL_VOCAB_SIZE-1})")
        
    except Exception as e:
        print(f"❌ 模型加载失败: {e}")
        raise
    
    return _session


# ---------- 模型加载 ----------
def load_models():
    """加载分词器和配置"""
    global _text_tokenizer, _filename_tokenizer, _categories, _config
    
    required_files = [
        (TEXT_TOKENIZER_PATH, "文本分词器"),
        (FILENAME_TOKENIZER_PATH, "文件名词典"),
        (CATEGORIES_PATH, "类别映射"),
        (CONFIG_PATH, "配置")
    ]
    
    for filepath, desc in required_files:
        if not Path(filepath).exists():
            raise FileNotFoundError(f"{desc}文件不存在: {filepath}")
    
    print("\n📚 加载模型组件...")
    print("-" * 40)
    
    print("加载文本分词器...", flush=True)
    _text_tokenizer = load_tokenizer_compressed(TEXT_TOKENIZER_PATH, MODEL_VOCAB_SIZE)
    print("  ✓")
    
    print("加载文件名词典...", flush=True)
    _filename_tokenizer = load_tokenizer_compressed(FILENAME_TOKENIZER_PATH, MODEL_VOCAB_SIZE)
    print("  ✓")
    
    print("加载类别映射...", end=' ', flush=True)
    with open(CATEGORIES_PATH, 'rb') as f:
        _categories = pickle.load(f)
    print("✓")
    
    print("加载配置...", end=' ', flush=True)
    with open(CONFIG_PATH, 'rb') as f:
        _config = pickle.load(f)
    print("✓")
    
    print("-" * 40)
    print(f"支持 {len(_categories)} 个类别: {', '.join(_categories)}")
    
    # 获取ONNX会话
    get_onnx_session()
    
    # 显示词表信息
    text_info = _text_tokenizer.get_vocab_info()
    filename_info = _filename_tokenizer.get_vocab_info()
    
    print(f"\n词表信息汇总:")
    print(f"  文本词表: {text_info['compressed_size']}/{text_info['original_size']} 个词 (保留率: {text_info['compressed_size']/text_info['original_size']*100:.1f}%)")
    print(f"  文件名词表: {filename_info['compressed_size']}/{filename_info['original_size']} 个词 (保留率: {filename_info['compressed_size']/filename_info['original_size']*100:.1f}%)")
    print(f"  模型支持最大词数: {MODEL_VOCAB_SIZE - 1}")
    
    # 显示OCR状态
    if OCR_AVAILABLE:
        print(f"\n✅ Tesseract OCR已启用（文本少于{OCR_MIN_TEXT_THRESHOLD}字符时自动识别图片文字）")
    else:
        print(f"\n⚠️ Tesseract OCR未安装（文本较少时可能影响分类效果）")
    print()


# ---------- OCR 图片文字识别 ----------
def extract_text_from_image(image_data):
    """从图片数据中提取中英文文字"""
    if not OCR_AVAILABLE:
        return ""
    
    try:
        # 从字节数据打开图片
        img = Image.open(BytesIO(image_data))
        
        # 配置Tesseract参数：中英文混合识别
        custom_config = r'--oem 3 --psm 6 -l chi_sim+eng'
        
        # 执行OCR识别
        text = pytesseract.image_to_string(img, config=custom_config)
        
        # 清理识别结果
        if text:
            text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    except Exception as e:
        # OCR失败时静默返回空字符串
        return ""


def extract_images_from_pptx(pptx_path):
    """从PPTX文件中提取所有图片并识别文字"""
    ocr_texts = []
    
    if not OCR_AVAILABLE:
        return ""
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            # PPTX中图片通常存储在 ppt/media/ 目录下
            image_extensions = ('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff', '.webp')
            
            for file_info in zf.filelist:
                if file_info.filename.startswith('ppt/media/') and \
                   file_info.filename.lower().endswith(image_extensions):
                    try:
                        # 读取图片数据
                        image_data = zf.read(file_info.filename)
                        
                        # 识别图片中的文字
                        recognized_text = extract_text_from_image(image_data)
                        if recognized_text and len(recognized_text) > 5:
                            ocr_texts.append(recognized_text)
                            
                    except Exception as e:
                        # 忽略单张图片的识别错误
                        continue
    
    except Exception as e:
        pass
    
    return " ".join(ocr_texts) if ocr_texts else ""


def extract_images_from_pptx_with_pptx_lib(pptx_path):
    """使用python-pptx库提取幻灯片中的图片（备用方法）"""
    ocr_texts = []
    
    if not OCR_AVAILABLE:
        return ""
    
    try:
        prs = Presentation(pptx_path)
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # 检查是否为图片形状
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # 获取图片
                        image = shape.image
                        image_bytes = image.blob
                        
                        # 识别图片中的文字
                        recognized_text = extract_text_from_image(image_bytes)
                        if recognized_text and len(recognized_text) > 5:
                            ocr_texts.append(recognized_text)
                    except Exception as e:
                        continue
                
                # 检查是否为包含图片的占位符
                elif hasattr(shape, 'has_text_frame') and not shape.has_text_frame:
                    try:
                        if hasattr(shape, 'image'):
                            image_bytes = shape.image.blob
                            recognized_text = extract_text_from_image(image_bytes)
                            if recognized_text and len(recognized_text) > 5:
                                ocr_texts.append(recognized_text)
                    except Exception:
                        pass
        
        del prs
        gc.collect()
        
    except Exception as e:
        pass
    
    return " ".join(ocr_texts) if ocr_texts else ""


def extract_ocr_from_pptx(pptx_path):
    """从PPTX文件中提取所有图片中的文字（结合两种方法）"""
    ocr_text = ""
    
    # 方法1：直接解压提取（更全面）
    ocr_text = extract_images_from_pptx(pptx_path)
    
    # 如果方法1没有提取到足够文字，尝试方法2
    if len(ocr_text) < 50:
        ocr_text2 = extract_images_from_pptx_with_pptx_lib(pptx_path)
        if ocr_text2:
            ocr_text = ocr_text + " " + ocr_text2 if ocr_text else ocr_text2
    
    return ocr_text


# ---------- DOCX 文本提取 ----------
def extract_text_from_docx(docx_path):
    """从DOCX文件中提取所有文本内容"""
    text_chunks = []
    total_text = ""
    doc = None
    
    if not DOCX_SUPPORT:
        print("  ⚠ 警告: python-docx未安装，无法解析docx文件")
        return ""
    
    try:
        doc = Document(docx_path)
        
        # 提取段落文本
        for paragraph in doc.paragraphs:
            if paragraph.text and paragraph.text.strip():
                text_chunks.append(paragraph.text.strip())
                if len(text_chunks) >= 20:
                    total_text += " ".join(text_chunks) + " "
                    text_chunks.clear()
        
        # 提取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        text_chunks.append(cell.text.strip())
                        if len(text_chunks) >= 20:
                            total_text += " ".join(text_chunks) + " "
                            text_chunks.clear()
        
        # 提取页眉页脚
        if hasattr(doc, 'sections'):
            for section in doc.sections:
                # 页眉
                if section.header and section.header.paragraphs:
                    for para in section.header.paragraphs:
                        if para.text and para.text.strip():
                            text_chunks.append(para.text.strip())
                # 页脚
                if section.footer and section.footer.paragraphs:
                    for para in section.footer.paragraphs:
                        if para.text and para.text.strip():
                            text_chunks.append(para.text.strip())
        
        # 添加剩余的文本块
        if text_chunks:
            total_text += " ".join(text_chunks)
            text_chunks.clear()
        
    except Exception as e:
        total_text = ""
    finally:
        if doc is not None:
            del doc
        text_chunks.clear()
        del text_chunks
        gc.collect()
    
    return total_text


# ---------- PPTX 嵌入DOCX提取 ----------
def extract_embedded_docx_from_pptx(pptx_path):
    """从PPTX文件中提取嵌入的DOCX文件并读取其内容"""
    embedded_texts = []
    
    try:
        # PPTX文件本质上是ZIP压缩包
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            # 查找嵌入的对象
            # 嵌入的文件通常在 ppt/embeddings/ 目录下
            for file_info in zf.filelist:
                if file_info.filename.startswith('ppt/embeddings/') and file_info.filename.endswith(('.docx', '.doc')):
                    try:
                        # 读取嵌入的docx文件内容
                        docx_data = zf.read(file_info.filename)
                        
                        # 使用临时文件或内存流处理docx
                        if DOCX_SUPPORT:
                            docx_stream = BytesIO(docx_data)
                            doc = Document(docx_stream)
                            
                            # 提取文本
                            for paragraph in doc.paragraphs:
                                if paragraph.text and paragraph.text.strip():
                                    embedded_texts.append(paragraph.text.strip())
                            
                            # 提取表格文本
                            for table in doc.tables:
                                for row in table.rows:
                                    for cell in row.cells:
                                        if cell.text and cell.text.strip():
                                            embedded_texts.append(cell.text.strip())
                            
                            del doc
                        else:
                            # 尝试简单提取（仅当python-docx不可用时）
                            # 简单方法：将docx作为zip读取，提取document.xml中的文本
                            try:
                                docx_zip = zipfile.ZipFile(BytesIO(docx_data))
                                if 'word/document.xml' in docx_zip.namelist():
                                    xml_content = docx_zip.read('word/document.xml').decode('utf-8', errors='ignore')
                                    # 简单正则提取文本
                                    text_matches = re.findall(r'>([^<]+)<', xml_content)
                                    for match in text_matches:
                                        if match.strip() and len(match.strip()) > 1:
                                            embedded_texts.append(match.strip())
                                docx_zip.close()
                            except:
                                pass
                        
                    except Exception as e:
                        # 忽略单个嵌入文件的解析错误
                        pass
    
    except Exception as e:
        pass
    
    return " ".join(embedded_texts) if embedded_texts else ""


# ---------- PPTX 文本提取（增强版，支持嵌入DOCX和OCR）----------
def extract_text_from_pptx(pptx_path, extract_embedded=True, use_ocr=True):
    """从PPTX文件中提取所有文本内容，包括嵌入的DOCX文件和图片OCR"""
    text_chunks = []
    total_text = ""
    prs = None
    embedded_text = ""
    ocr_text = ""
    
    try:
        # 提取嵌入的docx文件内容
        if extract_embedded:
            embedded_text = extract_embedded_docx_from_pptx(pptx_path)
            if embedded_text:
                text_chunks.append(embedded_text)
        
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    text_chunks.append(slide.shapes.title.text.strip())
                    if len(text_chunks) >= 10:
                        total_text += " ".join(text_chunks) + " "
                        text_chunks.clear()
            except:
                pass
            
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        text_chunks.append(shape.text.strip())
                        if len(text_chunks) >= 10:
                            total_text += " ".join(text_chunks) + " "
                            text_chunks.clear()
                    
                    # 提取表格中的文本
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text and cell.text.strip():
                                    text_chunks.append(cell.text.strip())
                                    if len(text_chunks) >= 10:
                                        total_text += " ".join(text_chunks) + " "
                                        text_chunks.clear()
                    
                    # 检查是否为嵌入的对象（备用方法）
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text_frame.text and shape.text_frame.text.strip():
                            text_chunks.append(shape.text_frame.text.strip())
                            
                except:
                    continue
            
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        text_chunks.append(notes_slide.notes_text_frame.text.strip())
                        if len(text_chunks) >= 10:
                            total_text += " ".join(text_chunks) + " "
                            text_chunks.clear()
            except:
                pass
            
            del slide
        
        if text_chunks:
            total_text += " ".join(text_chunks)
            text_chunks.clear()
        
    except Exception as e:
        total_text = ""
    finally:
        if prs is not None:
            del prs
        text_chunks.clear()
        del text_chunks
        gc.collect()
    
    # 如果提取的文本较少，启用OCR识别图片文字
    if use_ocr and OCR_AVAILABLE and len(total_text.strip()) < OCR_MIN_TEXT_THRESHOLD:
        ocr_text = extract_ocr_from_pptx(pptx_path)
        if ocr_text:
            # 将OCR识别的文字追加到总文本中
            if total_text:
                total_text = total_text + " " + ocr_text
            else:
                total_text = ocr_text
    
    return total_text, ocr_text if use_ocr and len(ocr_text) > 0 else ""


# ---------- 统一文件处理入口 ----------
def extract_text_from_file(filepath, use_ocr=True):
    """根据文件类型提取文本内容，返回(文本, OCR是否使用, OCR文本)"""
    filepath = Path(filepath)
    suffix = filepath.suffix.lower()
    
    if suffix in ['.pptx', '.ppt']:
        text, ocr_text = extract_text_from_pptx(str(filepath), use_ocr=use_ocr)
        return text, bool(ocr_text and len(ocr_text) > 0), ocr_text
    elif suffix == '.docx':
        return extract_text_from_docx(str(filepath)), False, ""
    else:
        return "", False, ""


# ---------- 文本预处理 ----------
def clean_text(text):
    """基础清洗"""
    if not text:
        return ""
    
    text = re.sub(r'http\S+|www\S+|https\S+', '', text)
    text = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9，。！？；：""''（）【】《》、 ]', '', text)
    text = re.sub(r'\s+', ' ', text)
    
    return text.strip()


def cut_words(text):
    """使用jieba进行分词"""
    if not text:
        return ""
    
    words = jieba.cut(text, HMM=False)
    filtered = []
    for w in words:
        if w and w.strip() and w not in STOPWORDS:
            filtered.append(w)
    
    if not filtered:
        words = jieba.cut(text, HMM=False)
        filtered = [w for w in words if w and w.strip()]
    
    if not filtered:
        return "空文本"
    
    result = " ".join(filtered)
    filtered.clear()
    del filtered
    del words
    
    return result


def extract_filename_features(filepath):
    """从文件路径中提取文件名特征"""
    filepath_obj = Path(filepath) if isinstance(filepath, str) else filepath
    filename = filepath_obj.stem
    
    cleaned = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', ' ', filename)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    
    words = jieba.cut(cleaned, HMM=False)
    filtered = []
    for w in words:
        if not w or not w.strip():
            continue
        if w not in STOPWORDS:
            if len(w) > 1 or w in MEANINGFUL_SINGLE_CHARS:
                filtered.append(w)
    
    if not filtered:
        words = jieba.cut(cleaned, HMM=False)
        filtered = [w for w in words if w and w.strip()]
    
    if not filtered:
        return cleaned
    
    result = " ".join(filtered)
    filtered.clear()
    del filtered
    del words
    
    return result


def pad_sequences(sequences, maxlen, padding='post', truncating='post'):
    """手动实现pad_sequences"""
    if not sequences:
        return np.zeros((0, maxlen), dtype=np.int32)
    
    result = np.zeros((len(sequences), maxlen), dtype=np.int32)
    
    for i, seq in enumerate(sequences):
        if truncating == 'post':
            seq = seq[:maxlen]
        else:
            seq = seq[-maxlen:] if len(seq) > maxlen else seq
        
        if padding == 'post':
            result[i, :len(seq)] = seq
        else:
            if len(seq) > 0:
                result[i, -len(seq):] = seq
    
    return result


# ---------- 核心预测函数（支持多种文件类型）----------
def predict_file(filepath, verbose=True):
    """预测单个文件（支持PPTX、PPT、DOCX）"""
    global _text_tokenizer, _filename_tokenizer, _categories, _config
    
    if _text_tokenizer is None:
        load_models()
    
    session = get_onnx_session()
    filepath = Path(filepath)
    
    if not filepath.exists():
        raise FileNotFoundError(f"文件不存在: {filepath}")
    
    suffix = filepath.suffix.lower()
    if suffix not in ['.pptx', '.ppt', '.docx']:
        raise ValueError(f"不支持的文件格式: {suffix}，仅支持 .pptx, .ppt, .docx")
    
    # 提取特征
    filename_feature = extract_filename_features(str(filepath))
    raw_text, ocr_used, ocr_text = extract_text_from_file(str(filepath), use_ocr=True)
    
    # 标记是否从嵌入docx中提取了内容
    embedded_used = False
    if suffix in ['.pptx', '.ppt'] and not raw_text.strip():
        # 如果PPTX本身没有文本，尝试更深入地提取嵌入docx
        embedded_text = extract_embedded_docx_from_pptx(str(filepath))
        if embedded_text:
            raw_text = embedded_text
            embedded_used = True
    
    # 记录原始文本来源信息
    text_sources = []
    if raw_text and len(raw_text.strip()) > 0:
        text_sources.append("幻灯片文本")
    if embedded_used:
        text_sources.append("嵌入DOCX")
    if ocr_used and ocr_text:
        text_sources.append("图片OCR")
    
    if raw_text:
        cleaned = clean_text(raw_text)
        del raw_text
        text_feature = cut_words(cleaned)
        del cleaned
        text_available = True
    else:
        text_feature = "空文本"
        text_available = False
        if verbose:
            print(f"  ⚠ 警告: 未提取到文本内容，仅使用文件名分类")
    
    text_len = len(text_feature) if text_feature else 0
    
    # 转换为序列（现在tokenizer会自动处理索引范围）
    text_seq = _text_tokenizer.texts_to_sequences([text_feature])
    text_pad = pad_sequences(
        text_seq, 
        maxlen=_config.get('max_sequence_length', 750),
        padding='post', 
        truncating='post'
    )
    del text_seq
    
    filename_seq = _filename_tokenizer.texts_to_sequences([filename_feature])
    filename_pad = pad_sequences(
        filename_seq, 
        maxlen=_config.get('max_filename_length', 30),
        padding='post', 
        truncating='post'
    )
    del filename_seq
    del text_feature, filename_feature
    
    # 确保数据类型正确
    text_pad = text_pad.astype(np.int32)
    filename_pad = filename_pad.astype(np.int32)
    
    # 调试：检查最大索引（应该都在有效范围内）
    if verbose:
        max_text_idx = text_pad.max()
        max_filename_idx = filename_pad.max()
        if max_text_idx >= MODEL_VOCAB_SIZE:
            print(f"  ⚠ 警告: 文本索引超出范围: {max_text_idx} >= {MODEL_VOCAB_SIZE}")
        if max_filename_idx >= MODEL_VOCAB_SIZE:
            print(f"  ⚠ 警告: 文件名索引超出范围: {max_filename_idx} >= {MODEL_VOCAB_SIZE}")
        print(f"  📊 索引范围: 文本[{text_pad.min()}-{max_text_idx}], 文件名[{filename_pad.min()}-{max_filename_idx}]")
    
    # 推理
    input_names = [inp.name for inp in session.get_inputs()]
    output_names = [out.name for out in session.get_outputs()]
    
    inputs = {
        input_names[0]: text_pad,
        input_names[1]: filename_pad
    }
    
    predictions = session.run(output_names, inputs)[0]
    del text_pad, filename_pad, inputs
    
    # 获取结果
    pred_idx = int(predictions[0].argmax())
    confidence = float(predictions[0][pred_idx])
    predicted_class = _categories[pred_idx]
    
    if verbose:
        all_probs = {cat: float(predictions[0][i]) for i, cat in enumerate(_categories)}
        sorted_probs = sorted(all_probs.items(), key=lambda x: x[1], reverse=True)
        
        print(f"\n{'='*50}")
        print(f"文件: {filepath.name}")
        print(f"类型: {suffix.upper()}")
        print(f"{'='*50}")
        print(f"📝 文本长度: {text_len} 词")
        print(f"📝 文本可用: {'是' if text_available else '否'}")
        if text_sources:
            print(f"📎 文本来源: {' + '.join(text_sources)}")
        if embedded_used:
            print(f"📎 来源: 从PPTX嵌入的DOCX文件中提取")
        if ocr_used and ocr_text:
            ocr_preview = ocr_text[:100] + "..." if len(ocr_text) > 100 else ocr_text
            print(f"🔍 OCR识别: {len(ocr_text)} 字符")
            if len(ocr_text) > 0:
                print(f"   OCR预览: {ocr_preview}")
        print(f"\n🎯 预测结果: {predicted_class}")
        print(f"📊 置信度: {confidence:.2%}")
        print(f"\n📈 详细分类概率:")
        for cat, prob in sorted_probs:
            bar_length = int(prob * 30)
            bar = "█" * bar_length if bar_length > 0 else "·"
            print(f"   {cat:4s}: {prob:6.2%} {bar}")
        
        del all_probs, sorted_probs
    else:
        source_info = []
        if embedded_used:
            source_info.append("embedded")
        if ocr_used and ocr_text:
            source_info.append("ocr")
        source_str = f" [{'+'.join(source_info)}]" if source_info else ""
        print(f"{filepath.name}{source_str}: {predicted_class} ({confidence:.2%})")
    
    del predictions
    gc.collect()
    
    return predicted_class, confidence, {
        'text_length': text_len,
        'text_available': text_available,
        'embedded_used': embedded_used,
        'ocr_used': ocr_used,
        'file_type': suffix,
        'text_sources': text_sources
    }


# 保持向后兼容
def predict_pptx(filepath, verbose=True):
    """预测单个PPTX文件（向后兼容）"""
    return predict_file(filepath, verbose)


# ---------- 交互式预测 ----------
def interactive_mode():
    """交互式预测模式"""
    print("\n" + "="*50)
    print("文件分类器 - 交互式预测模式")
    print("支持格式: PPTX, PPT, DOCX")
    if DOCX_SUPPORT:
        print("✅ python-docx已安装，支持完整DOCX解析")
    else:
        print("⚠️ python-docx未安装，DOCX支持受限（仅简单文本提取）")
    if OCR_AVAILABLE:
        print(f"✅ Tesseract OCR已启用（文本少于{OCR_MIN_TEXT_THRESHOLD}字符时自动识别图片文字）")
    else:
        print("⚠️ Tesseract OCR未安装（文本较少时可能影响分类效果）")
    print("="*50)
    
    # 显示词表信息
    text_info = _text_tokenizer.get_vocab_info()
    filename_info = _filename_tokenizer.get_vocab_info()
    
    print(f"文本词表: {text_info['compressed_size']} 个词 (原始: {text_info['original_size']})")
    print(f"文件名词表: {filename_info['compressed_size']} 个词 (原始: {filename_info['original_size']})")
    print(f"模型词汇表限制: {MODEL_VOCAB_SIZE} (有效索引: 1-{MODEL_VOCAB_SIZE-1})")
    print(f"OOV Token: {text_info['oov_token']}")
    print("="*50)
    print("输入文件路径进行预测（.pptx, .ppt, .docx），输入 'quit' 或 'exit' 退出")
    print("-"*50)
    
    while True:
        try:
            filepath = input("\n请输入文件路径: ").strip().strip('"').strip("'")
        except (KeyboardInterrupt, EOFError):
            print("\n再见！")
            break
        
        if filepath.lower() in ['quit', 'exit', 'q', 'bye']:
            print("再见！")
            break
        
        if not filepath:
            continue
        
        if not Path(filepath).exists():
            print(f"错误: 文件不存在 '{filepath}'")
            continue
        
        suffix = Path(filepath).suffix.lower()
        if suffix not in ['.pptx', '.ppt', '.docx']:
            print(f"错误: 不支持的文件格式 '{suffix}'，请使用 .pptx, .ppt 或 .docx 文件")
            continue
        
        try:
            predict_file(filepath, verbose=True)
            gc.collect()
        except Exception as e:
            print(f"预测失败: {e}")
            import traceback
            traceback.print_exc()


# ---------- 批量预测（支持多种文件类型）----------
def predict_batch(directory_path, recursive=False):
    """批量预测，支持PPTX和DOCX文件"""
    directory = Path(directory_path)
    
    if not directory.exists():
        raise FileNotFoundError(f"目录不存在: {directory}")
    
    # 查找所有支持的文件
    supported_extensions = ['*.pptx', '*.PPTX', '*.ppt', '*.PPT', '*.docx', '*.DOCX']
    pptx_files = []
    
    if recursive:
        for ext in supported_extensions:
            pptx_files.extend(directory.rglob(ext))
    else:
        for ext in supported_extensions:
            pptx_files.extend(directory.glob(ext))
    
    pptx_files = list(set(pptx_files))
    
    if not pptx_files:
        print(f"未找到支持的文件（.pptx, .ppt, .docx）: {directory}")
        return
    
    print(f"\n找到 {len(pptx_files)} 个文件，开始批量预测...")
    print("="*60)
    
    output_csv = directory / f"prediction_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
    import csv
    
    success_count = 0
    class_counts = {}
    type_counts = {}
    ocr_used_count = 0
    
    with open(output_csv, 'w', newline='', encoding='utf-8-sig') as f:
        writer = csv.DictWriter(f, fieldnames=['file', 'file_type', 'predicted_class', 'confidence', 
                                                'text_available', 'embedded_used', 'ocr_used', 'text_sources'])
        writer.writeheader()
        
        for i, filepath in enumerate(pptx_files, 1):
            suffix = filepath.suffix.lower()
            print(f"\n[{i}/{len(pptx_files)}] 处理中: {filepath.name[:50]}", end=' ')
            sys.stdout.flush()
            
            try:
                pred_class, confidence, info = predict_file(str(filepath), verbose=False)
                writer.writerow({
                    'file': str(filepath),
                    'file_type': suffix,
                    'predicted_class': pred_class,
                    'confidence': f"{confidence:.4f}",
                    'text_available': str(info.get('text_available', False)),
                    'embedded_used': str(info.get('embedded_used', False)),
                    'ocr_used': str(info.get('ocr_used', False)),
                    'text_sources': '+'.join(info.get('text_sources', []))
                })
                f.flush()
                success_count += 1
                class_counts[pred_class] = class_counts.get(pred_class, 0) + 1
                type_counts[suffix] = type_counts.get(suffix, 0) + 1
                if info.get('ocr_used'):
                    ocr_used_count += 1
                
                source_info = []
                if info.get('embedded_used'):
                    source_info.append("embedded")
                if info.get('ocr_used'):
                    source_info.append("ocr")
                source_str = f" [{'+'.join(source_info)}]" if source_info else ""
                print(f"✓ → {pred_class} ({confidence:.1%}){source_str}")
            except Exception as e:
                print(f"✗ → 错误: {str(e)[:50]}")
                writer.writerow({
                    'file': str(filepath),
                    'file_type': suffix,
                    'predicted_class': 'ERROR',
                    'confidence': '0',
                    'text_available': 'False',
                    'embedded_used': 'False',
                    'ocr_used': 'False',
                    'text_sources': ''
                })
                f.flush()
            
            gc.collect()
    
    print("\n" + "="*60)
    print("批量预测汇总")
    print("="*60)
    print(f"总计: {len(pptx_files)} 个文件")
    print(f"成功: {success_count} 个")
    print(f"失败: {len(pptx_files) - success_count} 个")
    
    if OCR_AVAILABLE and ocr_used_count > 0:
        print(f"🔍 OCR启用: {ocr_used_count} 个文件使用了图片文字识别")
    
    if type_counts:
        print("\n文件类型分布:")
        for ftype, count in sorted(type_counts.items(), key=lambda x: x[1], reverse=True):
            print(f"  {ftype}: {count} 个")
    
    if class_counts:
        print("\n类别分布:")
        for cat, count in sorted(class_counts.items(), key=lambda x: x[1], reverse=True):
            percentage = count / success_count * 100 if success_count > 0 else 0
            print(f"  {cat}: {count} 个 ({percentage:.1f}%)")
    
    print(f"\n📄 结果已导出到: {output_csv}")


# ---------- 命令行入口 ----------
def main():
    print("正在初始化ONNX推理引擎...")
    print(f"模型词汇表限制: {MODEL_VOCAB_SIZE} (有效索引: 1-{MODEL_VOCAB_SIZE-1})")
    print(f"DOCX支持: {'✅ 已启用' if DOCX_SUPPORT else '❌ 未安装python-docx'}")
    print(f"OCR支持: {'✅ 已启用（文本<{}字符时自动识别）'.format(OCR_MIN_TEXT_THRESHOLD) if OCR_AVAILABLE else '❌ 未安装pytesseract'}")
    print("-" * 40)
    
    try:
        load_models()
    except Exception as e:
        print(f"\n❌ 错误: 模型加载失败 - {e}")
        print("\n请确保以下文件存在于当前目录:")
        for f in [ONNX_MODEL_PATH, TEXT_TOKENIZER_PATH, FILENAME_TOKENIZER_PATH, CATEGORIES_PATH, CONFIG_PATH]:
            exists = "✓" if Path(f).exists() else "✗"
            print(f"  {exists} {f}")
        return 1
    
    parser = argparse.ArgumentParser(description='文件分类器预测工具（支持PPTX、PPT、DOCX，支持OCR图片文字识别）')
    parser.add_argument('input', nargs='?', help='文件路径或目录路径（支持.pptx, .ppt, .docx）')
    parser.add_argument('--batch', '-b', action='store_true', help='批量预测模式')
    parser.add_argument('--recursive', '-r', action='store_true', help='递归搜索子目录')
    parser.add_argument('--interactive', '-i', action='store_true', help='交互式预测模式')
    parser.add_argument('--no-ocr', action='store_true', help='禁用OCR图片文字识别')
    
    args = parser.parse_args()
    
    # 如果禁用OCR，设置全局标志
    if args.no_ocr:
        global OCR_AVAILABLE
        OCR_AVAILABLE = False
        print("⚠️ OCR已通过命令行参数禁用")
    
    try:
        if args.interactive:
            interactive_mode()
        elif args.batch or (args.input and Path(args.input).is_dir()):
            input_path = args.input or "."
            predict_batch(input_path, recursive=args.recursive)
        elif args.input:
            predict_file(args.input, verbose=True)
        else:
            parser.print_help()
    except KeyboardInterrupt:
        print("\n\n用户中断")
    except Exception as e:
        print(f"\n❌ 错误: {e}")
        import traceback
        traceback.print_exc()
        return 1
    
    return 0


if __name__ == "__main__":
    sys.exit(main())